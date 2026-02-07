#!/usr/bin/env python3
from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    import cairosvg  # type: ignore
except Exception:
    cairosvg = None

try:
    import requests  # type: ignore
except Exception:
    requests = None

PX_PER_INCH = 96.0
EMU_PER_INCH = 914400.0
DEFAULT_BASE_SLIDE_WIDTH_IN = 13.333
DEFAULT_BASE_SLIDE_HEIGHT_IN = 7.5
DEFAULT_MAX_SLIDE_WIDTH_IN = 56.0
DEFAULT_MAX_SLIDE_HEIGHT_IN = 31.5
# Keep a small outer gutter so diagrams visually use almost the full slide.
SLIDE_MARGIN_IN = 0.12

TOP = 0
LEFT = 1
BOTTOM = 2
RIGHT = 3
SIDE_TOKEN_MAP = {"T": TOP, "L": LEFT, "B": BOTTOM, "R": RIGHT}

SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "circle": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "parallelogramAlt": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "trapezoidAlt": MSO_SHAPE.NON_ISOSCELES_TRAPEZOID,
    "cylinder": MSO_SHAPE.CAN,
    "subroutine": MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
    "hexagon": MSO_SHAPE.HEXAGON,
    "cloud": MSO_SHAPE.CLOUD,
    "explosion": MSO_SHAPE.EXPLOSION1,
    "lightningBolt": MSO_SHAPE.LIGHTNING_BOLT,
    "notchedRightArrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
    "hourglass": MSO_SHAPE.FLOWCHART_COLLATE,
    "hCylinder": MSO_SHAPE.FLOWCHART_DIRECT_ACCESS_STORAGE,
    "curvedTrapezoid": MSO_SHAPE.FLOWCHART_DISPLAY,
    "forkBar": MSO_SHAPE.RECTANGLE,
    "windowPane": MSO_SHAPE.RECTANGLE,
    "filledCircle": MSO_SHAPE.OVAL,
    "smallCircle": MSO_SHAPE.OVAL,
    "framedCircle": MSO_SHAPE.OVAL,
    "linedDocument": MSO_SHAPE.FLOWCHART_DOCUMENT,
    "linedRect": MSO_SHAPE.RECTANGLE,
    "wave": MSO_SHAPE.WAVE,
    "stackedRect": MSO_SHAPE.RECTANGLE,
    "framedRect": MSO_SHAPE.RECTANGLE,
    "braceLeft": MSO_SHAPE.LEFT_BRACE,
    "braceRight": MSO_SHAPE.RIGHT_BRACE,
    "bracePair": MSO_SHAPE.DOUBLE_BRACE,
    "card": MSO_SHAPE.FLOWCHART_CARD,
    "delay": MSO_SHAPE.FLOWCHART_DELAY,
    "internalStorage": MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE,
    "document": MSO_SHAPE.FLOWCHART_DOCUMENT,
    "multiDocument": MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "rightTriangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "chevron": MSO_SHAPE.CHEVRON,
    "plaqueTabs": MSO_SHAPE.PLAQUE_TABS,
    "pentagon": MSO_SHAPE.PENTAGON,
    "decagon": MSO_SHAPE.DECAGON,
    "foldedCorner": MSO_SHAPE.FOLDED_CORNER,
    "donut": MSO_SHAPE.DONUT,
    "summingJunction": MSO_SHAPE.FLOWCHART_SUMMING_JUNCTION,
}

ARCH_THEME_BY_PROVIDER = {
    "aws": {"fill": "FFF7ED", "stroke": "FB923C", "text": "7C2D12"},
    "azure": {"fill": "EFF6FF", "stroke": "0EA5E9", "text": "0C4A6E"},
    "gcp": {"fill": "ECFDF5", "stroke": "22C55E", "text": "14532D"},
    "openai": {"fill": "EEF2FF", "stroke": "818CF8", "text": "312E81"},
    "terraform": {"fill": "F5F3FF", "stroke": "8B5CF6", "text": "4C1D95"},
}

ICONIFY_DIRECT_PREFIXES = {
    "simple-icons",
    "mdi",
    "tabler",
    "logos",
    "lucide",
    "ph",
    "ri",
    "bi",
    "material-symbols",
}

AWS_SERVICE_ICON_MAP = {
    "s3": "simple-icons/amazons3",
    "rds": "simple-icons/amazonrds",
    "postgres": "simple-icons/postgresql",
    "postgresql": "simple-icons/postgresql",
    "lambda": "simple-icons/awslambda",
    "apigateway": "simple-icons/amazonapigateway",
    "api-gateway": "simple-icons/amazonapigateway",
    "ecs": "simple-icons/amazonecs",
    "ecr": "simple-icons/amazonecr",
    "sqs": "simple-icons/amazonsqs",
    "sns": "simple-icons/amazonsns",
    "dynamodb": "simple-icons/amazondynamodb",
    "eventbridge": "simple-icons/amazoneventbridge",
    "cloudwatch": "simple-icons/amazoncloudwatch",
    "kms": "simple-icons/amazonkms",
}


def px_to_in(value: float) -> float:
    return value / PX_PER_INCH


def emu_to_in(value: float) -> float:
    return value / EMU_PER_INCH


def to_rgb(value: str) -> RGBColor:
    text = (value or "000000").strip().lstrip("#")
    if len(text) != 6:
        text = "000000"
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def blend_with_white(value: str, ratio: float = 0.72) -> str:
    text = (value or "000000").strip().lstrip("#")
    if len(text) != 6:
        text = "000000"

    ratio = max(0.0, min(1.0, ratio))
    r = int(text[0:2], 16)
    g = int(text[2:4], 16)
    b = int(text[4:6], 16)

    wr = int(r * (1.0 - ratio) + 255 * ratio)
    wg = int(g * (1.0 - ratio) + 255 * ratio)
    wb = int(b * (1.0 - ratio) + 255 * ratio)
    return f"{wr:02X}{wg:02X}{wb:02X}"


def _as_slug(value: str) -> str:
    text = re.sub(r"[^a-z0-9]+", "-", value.strip().lower())
    return text.strip("-")


def _provider_from_text(text: str) -> str | None:
    lowered = text.lower()
    if "openai" in lowered:
        return "openai"
    if "azure" in lowered or "microsoft" in lowered:
        return "azure"
    if "aws" in lowered or "amazon" in lowered:
        return "aws"
    if "gcp" in lowered or "google cloud" in lowered:
        return "gcp"
    if "terraform" in lowered:
        return "terraform"
    return None


def infer_arch_theme(node: dict[str, Any]) -> dict[str, str] | None:
    icon_raw = str(node.get("icon", "")).strip()
    if not icon_raw:
        return None
    label = str(node.get("label", ""))
    node_id = str(node.get("id", ""))
    provider = _provider_from_text(f"{icon_raw} {label} {node_id}")
    if not provider:
        return None
    return ARCH_THEME_BY_PROVIDER.get(provider)


def normalize_iconify_key(raw: str) -> str:
    text = raw.strip().lower()
    if not text:
        return ""
    text = text.replace(" ", "-").replace("_", "-")
    if ":" in text:
        prefix, name = text.split(":", 1)
        return f"{prefix}/{name}"
    return text


def iconify_candidates(icon_raw: str, label: str, node_id: str) -> list[str]:
    out: list[str] = []
    raw = (icon_raw or "").strip().lower()
    label_slug = _as_slug(label)
    id_slug = _as_slug(node_id)

    if raw:
        normalized = normalize_iconify_key(raw)
        if "/" in normalized:
            prefix = normalized.split("/", 1)[0]
            if prefix in ICONIFY_DIRECT_PREFIXES:
                out.append(normalized)

        if raw.startswith("aws:") or normalized.startswith("aws/"):
            service = normalize_iconify_key(raw).split("/", 1)[1] if "/" in normalize_iconify_key(raw) else ""
            service = service.replace("amazon-", "").replace("aws-", "")
            if service in AWS_SERVICE_ICON_MAP:
                out.append(AWS_SERVICE_ICON_MAP[service])
            out.append(f"simple-icons/amazon{service}")
            out.append("simple-icons/amazonwebservices")
        elif raw.startswith("azure:") or normalized.startswith("azure/"):
            service = normalize_iconify_key(raw).split("/", 1)[1] if "/" in normalize_iconify_key(raw) else ""
            out.append(f"simple-icons/azure{service}")
            out.append("simple-icons/microsoftazure")
        elif raw.startswith("gcp:") or normalized.startswith("gcp/"):
            service = normalize_iconify_key(raw).split("/", 1)[1] if "/" in normalize_iconify_key(raw) else ""
            out.append(f"simple-icons/googlecloud{service}")
            out.append("simple-icons/googlecloud")
        elif raw.startswith("openai:") or normalized.startswith("openai/"):
            out.append("simple-icons/openai")
        elif raw in {"aws", "amazon"}:
            out.append("simple-icons/amazonwebservices")
        elif raw in {"azure", "microsoftazure"}:
            out.append("simple-icons/microsoftazure")
        elif raw in {"openai"}:
            out.append("simple-icons/openai")
        elif raw in {"terraform"}:
            out.append("simple-icons/terraform")

    inferred_provider = _provider_from_text(f"{raw} {label} {node_id}")
    if inferred_provider == "aws":
        out.extend(
            [
                "simple-icons/amazonwebservices",
                "simple-icons/amazons3" if "s3" in f"{label_slug}-{id_slug}" else "",
                "simple-icons/postgresql" if "postgres" in f"{label_slug}-{id_slug}" else "",
            ]
        )
    elif inferred_provider == "azure":
        out.append("simple-icons/microsoftazure")
    elif inferred_provider == "openai":
        out.append("simple-icons/openai")
    elif inferred_provider == "terraform":
        out.append("simple-icons/terraform")

    if "s3" in f"{label_slug}-{id_slug}":
        out.append("simple-icons/amazons3")
    if "postgres" in f"{label_slug}-{id_slug}":
        out.append("simple-icons/postgresql")

    compact: list[str] = []
    seen: set[str] = set()
    for item in out:
        key = item.strip().lower()
        if not key or key in seen:
            continue
        compact.append(key)
        seen.add(key)
    return compact


def icon_cache_dir() -> Path:
    return Path.home() / ".cache" / "mmd2pptx" / "icons"


def resolve_icon_png(icon_raw: str, label: str, node_id: str, size_px: int = 96) -> Path | None:
    if not icon_raw or not icon_raw.strip():
        return None
    if requests is None or cairosvg is None:
        return None

    cache_dir = icon_cache_dir()
    cache_dir.mkdir(parents=True, exist_ok=True)
    timeout = (4, 8)

    for icon_key in iconify_candidates(icon_raw, label, node_id):
        digest = hashlib.sha1(f"{icon_key}:{size_px}".encode("utf-8")).hexdigest()[:16]
        png_path = cache_dir / f"{digest}.png"
        if png_path.exists() and png_path.stat().st_size > 0:
            return png_path

        url = f"https://api.iconify.design/{icon_key}.svg"
        try:
            response = requests.get(url, timeout=timeout, headers={"User-Agent": "mmd2pptx/0.1"})
            if response.status_code != 200:
                continue
            svg_bytes = response.content
            if b"<svg" not in svg_bytes[:256]:
                continue
            cairosvg.svg2png(bytestring=svg_bytes, write_to=str(png_path), output_width=size_px, output_height=size_px)
            if png_path.exists() and png_path.stat().st_size > 0:
                return png_path
        except Exception:
            continue

    return None


def parse_slide_size(slide_size: str | None) -> tuple[float, float, float, float]:
    text = (slide_size or "16:9").strip().lower()

    if text in {"16:9", "wide", "widescreen"}:
        # Respect requested preset size exactly.
        return 13.333, 7.5, 13.333, 7.5

    if text in {"4:3", "standard"}:
        # Respect requested preset size exactly.
        return 10.0, 7.5, 10.0, 7.5

    if "x" in text:
        try:
            w_text, h_text = text.split("x", 1)
            base_w = float(w_text.strip())
            base_h = float(h_text.strip())
            if base_w <= 0 or base_h <= 0:
                raise ValueError("non-positive slide size")
        except Exception as exc:
            raise ValueError(f"invalid --slide-size '{slide_size}' (expected 16:9, 4:3, or WxH)") from exc

        # Respect requested custom size exactly.
        return base_w, base_h, base_w, base_h

    raise ValueError(f"invalid --slide-size '{slide_size}' (expected 16:9, 4:3, or WxH)")


def resolve_slide_size_and_scale(
    ir: dict[str, Any],
    *,
    base_slide_w: float,
    base_slide_h: float,
    max_slide_w: float,
    max_slide_h: float,
) -> tuple[float, float, float, float, float]:
    bounds = ir.get("bounds", {})
    width = max(float(bounds.get("width", 1.0)), 1.0)
    height = max(float(bounds.get("height", 1.0)), 1.0)

    width_in = px_to_in(width)
    height_in = px_to_in(height)

    # Keep requested slide aspect ratio strictly, and only scale both axes uniformly.
    required_w = width_in + SLIDE_MARGIN_IN * 2
    required_h = height_in + SLIDE_MARGIN_IN * 2

    growth_w = required_w / max(base_slide_w, 1e-6)
    growth_h = required_h / max(base_slide_h, 1e-6)
    growth = max(1.0, growth_w, growth_h)

    max_growth = min(max_slide_w / max(base_slide_w, 1e-6), max_slide_h / max(base_slide_h, 1e-6))
    growth = min(growth, max_growth)

    slide_w = base_slide_w * growth
    slide_h = base_slide_h * growth

    available_w = slide_w - SLIDE_MARGIN_IN * 2
    available_h = slide_h - SLIDE_MARGIN_IN * 2
    scale = min(available_w / width_in, available_h / height_in)

    content_w = width_in * scale
    content_h = height_in * scale
    min_x_in = px_to_in(float(bounds.get("minX", 0.0)))
    min_y_in = px_to_in(float(bounds.get("minY", 0.0)))

    # Center the diagram inside the slide viewport after scaling.
    offset_x = (slide_w - content_w) / 2.0 - min_x_in * scale
    offset_y = (slide_h - content_h) / 2.0 - min_y_in * scale
    return slide_w, slide_h, scale, offset_x, offset_y


def transform(x: float, y: float, scale: float, offset_x: float, offset_y: float) -> tuple[float, float]:
    return px_to_in(x) * scale + offset_x, px_to_in(y) * scale + offset_y


def side_point(node: dict[str, Any], side: int) -> tuple[float, float]:
    x = float(node["x"])
    y = float(node["y"])
    w = float(node["width"])
    h = float(node["height"])

    if side == TOP:
        return x + w / 2.0, y
    if side == LEFT:
        return x, y + h / 2.0
    if side == BOTTOM:
        return x + w / 2.0, y + h
    return x + w, y + h / 2.0


def side_from_token(value: Any) -> int | None:
    token = str(value).strip().upper()
    return SIDE_TOKEN_MAP.get(token)


def side_anchor_point(node: dict[str, Any], side: int, along_offset: float = 0.0) -> tuple[float, float]:
    x = float(node["x"])
    y = float(node["y"])
    w = float(node["width"])
    h = float(node["height"])
    margin = min(12.0, max(4.0, min(w, h) * 0.22))

    if side == TOP:
        ax = clampf(x + w / 2.0 + along_offset, x + margin, x + w - margin)
        return ax, y
    if side == BOTTOM:
        ax = clampf(x + w / 2.0 + along_offset, x + margin, x + w - margin)
        return ax, y + h
    if side == LEFT:
        ay = clampf(y + h / 2.0 + along_offset, y + margin, y + h - margin)
        return x, ay

    ay = clampf(y + h / 2.0 + along_offset, y + margin, y + h - margin)
    return x + w, ay


def side_normal(side: int) -> tuple[float, float]:
    if side == TOP:
        return 0.0, -1.0
    if side == LEFT:
        return -1.0, 0.0
    if side == BOTTOM:
        return 0.0, 1.0
    return 1.0, 0.0


def side_faces_vector(side: int, vx: float, vy: float) -> bool:
    nx, ny = side_normal(side)
    return nx * vx + ny * vy > 1e-6


def connection_cost(src: dict[str, Any], dst: dict[str, Any], src_side: int, dst_side: int) -> float:
    sx, sy = side_point(src, src_side)
    tx, ty = side_point(dst, dst_side)

    vx = tx - sx
    vy = ty - sy
    dist = (vx * vx + vy * vy) ** 0.5

    if dist < 1e-6:
        return 0.0

    src_nx, src_ny = side_normal(src_side)
    dst_nx, dst_ny = side_normal(dst_side)

    src_dot = src_nx * vx + src_ny * vy
    dst_dot = dst_nx * (-vx) + dst_ny * (-vy)

    # Keep distance minimization as the primary goal.
    # Apply only a moderate penalty when a side points opposite to the edge direction
    # to avoid lines that immediately dive into the source/target shape.
    penalty = 0.0
    if src_dot <= 0:
        penalty += 42.0 + dist * 0.42
    if dst_dot <= 0:
        penalty += 42.0 + dist * 0.42

    # Slightly prefer axis-consistent pairings when costs are near-equal.
    if (src_side in {TOP, BOTTOM}) != (dst_side in {TOP, BOTTOM}):
        penalty += 5.0

    return dist + penalty


def choose_connection_sides(
    src: dict[str, Any],
    dst: dict[str, Any],
    *,
    avoid_exact_pair: tuple[int, int] | None = None,
) -> tuple[int, int]:
    sides = (TOP, LEFT, BOTTOM, RIGHT)
    candidates = [(s, t) for s in sides for t in sides]

    best_src = RIGHT
    best_dst = LEFT
    best_cost = float("inf")
    for src_side, dst_side in candidates:
        cost = connection_cost(src, dst, src_side, dst_side)
        if avoid_exact_pair is not None and (src_side, dst_side) == avoid_exact_pair:
            cost += 2200.0
        if cost < best_cost:
            best_cost = cost
            best_src = src_side
            best_dst = dst_side

    return best_src, best_dst


def choose_connection_sides_with_hints(
    src: dict[str, Any],
    dst: dict[str, Any],
    *,
    hinted_src_side: int | None,
    hinted_dst_side: int | None,
    strict_hints: bool = False,
) -> tuple[int, int]:
    if strict_hints:
        if hinted_src_side is not None and hinted_dst_side is not None:
            return hinted_src_side, hinted_dst_side
        if hinted_src_side is not None:
            best_dst = min((TOP, LEFT, BOTTOM, RIGHT), key=lambda side: connection_cost(src, dst, hinted_src_side, side))
            return hinted_src_side, best_dst
        if hinted_dst_side is not None:
            best_src = min((TOP, LEFT, BOTTOM, RIGHT), key=lambda side: connection_cost(src, dst, side, hinted_dst_side))
            return best_src, hinted_dst_side

    auto_src, auto_dst = choose_connection_sides(src, dst)
    auto_cost = connection_cost(src, dst, auto_src, auto_dst)

    # Accept hinted sides only when not materially worse than shortest route.
    # This keeps diagrams readable when author-provided sides conflict with the
    # final layout.
    tolerance = max(10.0, auto_cost * 0.20)

    if hinted_src_side is not None and hinted_dst_side is not None:
        hinted_cost = connection_cost(src, dst, hinted_src_side, hinted_dst_side)
        if hinted_cost <= auto_cost + tolerance:
            return hinted_src_side, hinted_dst_side
        return auto_src, auto_dst

    if hinted_src_side is not None:
        best_dst = min((TOP, LEFT, BOTTOM, RIGHT), key=lambda side: connection_cost(src, dst, hinted_src_side, side))
        hinted_cost = connection_cost(src, dst, hinted_src_side, best_dst)
        if hinted_cost <= auto_cost + tolerance:
            return hinted_src_side, best_dst
        return auto_src, auto_dst

    if hinted_dst_side is not None:
        best_src = min((TOP, LEFT, BOTTOM, RIGHT), key=lambda side: connection_cost(src, dst, side, hinted_dst_side))
        hinted_cost = connection_cost(src, dst, best_src, hinted_dst_side)
        if hinted_cost <= auto_cost + tolerance:
            return best_src, hinted_dst_side
        return auto_src, auto_dst

    return auto_src, auto_dst


def side_sort_axis(side: int, from_node: dict[str, Any], to_node: dict[str, Any], *, source_side: bool) -> float:
    src = from_node if source_side else to_node
    dst = to_node if source_side else from_node
    dst_cx = float(dst["x"]) + float(dst["width"]) / 2.0
    dst_cy = float(dst["y"]) + float(dst["height"]) / 2.0
    src_cx = float(src["x"]) + float(src["width"]) / 2.0
    src_cy = float(src["y"]) + float(src["height"]) / 2.0

    if side in {TOP, BOTTOM}:
        return dst_cx + (dst_cy - src_cy) * 0.03
    return dst_cy + (dst_cx - src_cx) * 0.03


def lane_offsets(count: int, span_px: float) -> list[float]:
    if count <= 1:
        return [0.0]

    usable_span = max(0.0, span_px - 24.0)
    if usable_span <= 1e-6:
        step = 0.0
    else:
        step = min(15.0, max(6.0, usable_span / max(1, count)))

    center = (count - 1) / 2.0
    return [(idx - center) * step for idx in range(count)]


def clampf(value: float, min_v: float, max_v: float) -> float:
    return max(min_v, min(max_v, value))


def choose_self_loop_geometry(
    node_box: tuple[float, float, float, float],
    *,
    slide_w: float,
    slide_h: float,
) -> tuple[int, int, list[tuple[float, float]], tuple[float, float]]:
    x, y, w, h = node_box
    left_space = x
    right_space = slide_w - (x + w)
    top_space = y
    bottom_space = slide_h - (y + h)

    candidates = [
        ("tr", right_space + top_space),
        ("br", right_space + bottom_space),
        ("bl", left_space + bottom_space),
        ("tl", left_space + top_space),
    ]
    mode = max(candidates, key=lambda item: item[1])[0]
    loop_x = max(0.26, min(1.25, w * 0.72))
    loop_y = max(0.26, min(1.25, h * 1.05))

    if mode == "tr":
        src_side, dst_side = RIGHT, TOP
        points = [
            (x + w, y + h * 0.58),
            (x + w + loop_x, y + h * 0.58),
            (x + w + loop_x, y - loop_y),
            (x + w * 0.58, y - loop_y),
            (x + w * 0.58, y),
        ]
        label_anchor = (x + w + loop_x * 0.52, y - loop_y * 0.56)
    elif mode == "br":
        src_side, dst_side = RIGHT, BOTTOM
        points = [
            (x + w, y + h * 0.42),
            (x + w + loop_x, y + h * 0.42),
            (x + w + loop_x, y + h + loop_y),
            (x + w * 0.58, y + h + loop_y),
            (x + w * 0.58, y + h),
        ]
        label_anchor = (x + w + loop_x * 0.52, y + h + loop_y * 0.56)
    elif mode == "bl":
        src_side, dst_side = LEFT, BOTTOM
        points = [
            (x, y + h * 0.42),
            (x - loop_x, y + h * 0.42),
            (x - loop_x, y + h + loop_y),
            (x + w * 0.42, y + h + loop_y),
            (x + w * 0.42, y + h),
        ]
        label_anchor = (x - loop_x * 0.52, y + h + loop_y * 0.56)
    else:
        src_side, dst_side = LEFT, TOP
        points = [
            (x, y + h * 0.58),
            (x - loop_x, y + h * 0.58),
            (x - loop_x, y - loop_y),
            (x + w * 0.42, y - loop_y),
            (x + w * 0.42, y),
        ]
        label_anchor = (x - loop_x * 0.52, y - loop_y * 0.56)

    clamped_points: list[tuple[float, float]] = []
    for px, py in points:
        clamped_points.append((clampf(px, 0.03, slide_w - 0.03), clampf(py, 0.03, slide_h - 0.03)))

    lx = clampf(label_anchor[0], 0.08, slide_w - 0.08)
    ly = clampf(label_anchor[1], 0.08, slide_h - 0.08)
    return src_side, dst_side, clamped_points, (lx, ly)


def marker_to_ooxml(marker: str | None) -> str | None:
    token = (marker or "").strip()
    if token in {"none", ""}:
        return None
    if token == "arrow":
        return "arrow"
    if token == "triangle":
        return "triangle"
    if token == "diamond":
        return "diamond"
    if token in {"openDiamond", "circle"}:
        return None
    return None


def set_edge_markers(connector: Any, style: dict[str, Any]) -> None:
    ln = connector._element.spPr.get_or_add_ln()
    for child in list(ln):
        if child.tag in {qn("a:headEnd"), qn("a:tailEnd")}:
            ln.remove(child)

    start_marker = marker_to_ooxml(style.get("startMarker"))
    end_marker = marker_to_ooxml(style.get("endMarker"))

    start_marker_raw = str(style.get("startMarker", "")).strip()
    end_marker_raw = str(style.get("endMarker", "")).strip()
    has_explicit_markers = start_marker_raw != "" or end_marker_raw != ""

    if start_marker is None and end_marker is None and not has_explicit_markers:
        arrow = str(style.get("arrow", "none"))
        if arrow in {"start", "both"}:
            start_marker = "triangle"
        if arrow in {"end", "both"}:
            end_marker = "triangle"

    if start_marker:
        head = OxmlElement("a:headEnd")
        head.set("type", start_marker)
        ln.append(head)

    if end_marker:
        tail = OxmlElement("a:tailEnd")
        tail.set("type", end_marker)
        ln.append(tail)


def add_endpoint_symbol(
    group: Any,
    marker: str,
    *,
    at_start: bool,
    sx: float,
    sy: float,
    dx: float,
    dy: float,
    color: str,
) -> None:
    marker_type = (marker or "").strip()
    if marker_type not in {"openDiamond", "circle"}:
        return

    vx = dx - sx
    vy = dy - sy
    length = (vx * vx + vy * vy) ** 0.5
    if length < 1e-6:
        return

    tx = vx / length
    ty = vy / length
    if at_start:
        anchor_x = sx
        anchor_y = sy
        ux = tx
        uy = ty
    else:
        anchor_x = dx
        anchor_y = dy
        ux = -tx
        uy = -ty

    size = 0.12 if marker_type == "openDiamond" else 0.11
    cx = anchor_x + ux * (size * 0.42)
    cy = anchor_y + uy * (size * 0.42)
    left = cx - size / 2.0
    top = cy - size / 2.0

    shape_kind = MSO_SHAPE.DIAMOND if marker_type == "openDiamond" else MSO_SHAPE.OVAL
    symbol = group.shapes.add_shape(shape_kind, Inches(left), Inches(top), Inches(size), Inches(size))
    symbol.fill.solid()
    symbol.fill.fore_color.rgb = RGBColor(255, 255, 255)
    symbol.line.color.rgb = to_rgb(color)
    symbol.line.width = Pt(1.0)


def apply_line_style(line: Any, style: dict[str, Any]) -> None:
    width = float(style.get("width", 1.3))
    line_style = style.get("lineStyle", "solid")

    if line_style == "invisible":
        width = 0.1
    elif line_style == "thick":
        width = max(width * 2.3, 2.8)
    elif line_style == "dotted":
        width = max(width, 1.8)

    line.width = Pt(width)
    line.color.rgb = to_rgb(style.get("color", "1E293B"))
    if line_style == "invisible":
        line.transparency = 1.0
        line.dash_style = MSO_LINE_DASH_STYLE.SOLID
        return

    if line_style == "dotted":
        line.dash_style = MSO_LINE_DASH_STYLE.SQUARE_DOT
    else:
        line.dash_style = MSO_LINE_DASH_STYLE.SOLID


def scaled_font_size(base: float, scale: float) -> float:
    # Keep text proportional when the diagram must be scaled down to fit slide max size.
    return max(4.0, min(48.0, base * scale))


def text_units(text: str) -> float:
    units = 0.0
    for ch in text:
        if ch.isspace():
            units += 0.35
            continue
        code = ord(ch)
        if code >= 0x3000:
            units += 1.75
            continue
        units += 1.0
    return units


def estimate_label_box(text: str, font_size_pt: float) -> tuple[float, float]:
    units = max(2.0, text_units(text))
    width_in = (units * font_size_pt * 0.72) / 72.0 + 0.20
    width_in = max(0.9, min(4.8, width_in))
    height_in = max(0.26, min(0.95, (font_size_pt * 1.35) / 72.0 + 0.12))
    return width_in, height_in


def estimate_wrapped_line_count(text: str, width_in: float, font_size_pt: float) -> int:
    lines = [line for line in text.replace("\r\n", "\n").split("\n") if line.strip()]
    if not lines:
        return 1

    safe_width = max(0.25, width_in)
    units_per_line = max(4.0, (safe_width * 72.0) / max(font_size_pt * 0.62, 1e-6))
    wrapped = 0
    for line in lines:
        units = max(1.0, text_units(line.strip()))
        wrapped += max(1, int(math.ceil(units / units_per_line)))
    return max(1, wrapped)


def compact_text_to_fit_box(text: str, width_in: float, height_in: float, font_size_pt: float) -> str:
    lines = [line.strip() for line in text.replace("\r\n", "\n").split("\n") if line.strip()]
    if not lines:
        return ""

    max_lines = max(1, int((max(0.16, height_in) * 72.0 - 0.08) / max(font_size_pt * 1.26, 1e-6)))
    out: list[str] = []
    used_all = True

    for line in lines:
        candidate = out + [line]
        if estimate_wrapped_line_count("\n".join(candidate), width_in, font_size_pt) <= max_lines:
            out.append(line)
            continue
        used_all = False
        break

    if not out:
        out = [lines[0]]
        used_all = len(lines) <= 1

    if not used_all and out:
        trimmed = out[-1].rstrip(".")
        if not trimmed.endswith("..."):
            trimmed = f"{trimmed}..."
        out[-1] = trimmed

    while out and estimate_wrapped_line_count("\n".join(out), width_in, font_size_pt) > max_lines:
        last = out[-1]
        if len(last) <= 4:
            if len(out) <= 1:
                break
            out.pop()
            out[-1] = f"{out[-1].rstrip('.')}..."
            continue
        core = last[:-4].rstrip(" /,;:+-")
        out[-1] = f"{core}..."

    return "\n".join(out)


def rect_overlap(ax: float, ay: float, aw: float, ah: float, bx: float, by: float, bw: float, bh: float) -> bool:
    return not (ax + aw <= bx or bx + bw <= ax or ay + ah <= by or by + bh <= ay)


def rect_intersection_area(
    ax: float,
    ay: float,
    aw: float,
    ah: float,
    bx: float,
    by: float,
    bw: float,
    bh: float,
) -> float:
    left = max(ax, bx)
    top = max(ay, by)
    right = min(ax + aw, bx + bw)
    bottom = min(ay + ah, by + bh)
    if right <= left or bottom <= top:
        return 0.0
    return (right - left) * (bottom - top)


def segment_intersection(
    a0: tuple[float, float],
    a1: tuple[float, float],
    b0: tuple[float, float],
    b1: tuple[float, float],
) -> tuple[float, float] | None:
    x1, y1 = a0
    x2, y2 = a1
    x3, y3 = b0
    x4, y4 = b1

    denom = (x1 - x2) * (y3 - y4) - (y1 - y2) * (x3 - x4)
    if abs(denom) < 1e-9:
        return None

    t = ((x1 - x3) * (y3 - y4) - (y1 - y3) * (x3 - x4)) / denom
    u = ((x1 - x3) * (y1 - y2) - (y1 - y3) * (x1 - x2)) / denom

    if not (0.02 <= t <= 0.98 and 0.02 <= u <= 0.98):
        return None

    return x1 + t * (x2 - x1), y1 + t * (y2 - y1)


def overlap_count_with_obstacles(
    cx: float,
    cy: float,
    bw: float,
    bh: float,
    obstacle_boxes: list[tuple[float, float, float, float]],
) -> int:
    lx = cx - bw / 2.0
    ly = cy - bh / 2.0
    count = 0
    for nx, ny, nw, nh in obstacle_boxes:
        if rect_overlap(lx, ly, bw, bh, nx, ny, nw, nh):
            count += 1
    return count


def overlap_area_with_obstacles(
    cx: float,
    cy: float,
    bw: float,
    bh: float,
    obstacle_boxes: list[tuple[float, float, float, float]],
) -> float:
    lx = cx - bw / 2.0
    ly = cy - bh / 2.0
    area = 0.0
    for nx, ny, nw, nh in obstacle_boxes:
        area += rect_intersection_area(lx, ly, bw, bh, nx, ny, nw, nh)
    return area


def overlap_count_with_labels(
    cx: float,
    cy: float,
    bw: float,
    bh: float,
    label_boxes: list[tuple[float, float, float, float]],
    margin: float = 0.08,
) -> int:
    lx = cx - bw / 2.0
    ly = cy - bh / 2.0
    count = 0
    for px, py, pw, ph in label_boxes:
        if rect_overlap(lx, ly, bw, bh, px - margin, py - margin, pw + margin * 2, ph + margin * 2):
            count += 1
    return count


def overlap_area_with_labels(
    cx: float,
    cy: float,
    bw: float,
    bh: float,
    label_boxes: list[tuple[float, float, float, float]],
    margin: float = 0.08,
) -> float:
    lx = cx - bw / 2.0
    ly = cy - bh / 2.0
    area = 0.0
    for px, py, pw, ph in label_boxes:
        area += rect_intersection_area(lx, ly, bw, bh, px - margin, py - margin, pw + margin * 2, ph + margin * 2)
    return area


def intersection_penalty(
    cx: float,
    cy: float,
    intersections: list[tuple[float, float]],
    *,
    min_dist: float = 0.40,
) -> float:
    if not intersections:
        return 0.0

    penalty = 0.0
    for ix, iy in intersections:
        dist = ((cx - ix) ** 2 + (cy - iy) ** 2) ** 0.5
        if dist < min_dist:
            penalty += (min_dist - dist) * 10000.0
    return penalty


def choose_label_center(
    sx: float,
    sy: float,
    dx: float,
    dy: float,
    bw: float,
    bh: float,
    *,
    intersections: list[tuple[float, float]],
    obstacle_boxes: list[tuple[float, float, float, float]],
    placed_labels: list[tuple[float, float, float, float]],
    slide_w: float,
    slide_h: float,
    lane_t_bias: float = 0.0,
) -> tuple[float, float]:
    vx = dx - sx
    vy = dy - sy
    length = (vx * vx + vy * vy) ** 0.5

    if length < 1e-6:
        tx, ty = 1.0, 0.0
        nx, ny = 0.0, -1.0
    else:
        tx, ty = vx / length, vy / length
        nx, ny = -ty, tx

    # Short edges need larger normal offsets; otherwise labels are forced onto nearby nodes.
    short_pressure = max(0.0, bw * 0.75 - length)
    max_normal = min(1.05, max(0.16, length * 0.18, bh * 0.95 + 0.06, short_pressure * 0.70))
    if not intersections:
        max_normal = min(max_normal, 0.18)

    if intersections:
        # For crossed edges, avoid the center first so both labels retreat away from the intersection.
        t_positions = [0.20, 0.80, 0.32, 0.68, 0.44, 0.56, 0.14, 0.86]
    else:
        if length < max(0.95, bw * 0.60):
            t_positions = [0.50, 0.40, 0.60, 0.32, 0.68, 0.24, 0.76, 0.16, 0.84]
        else:
            t_positions = [0.50, 0.42, 0.58, 0.34, 0.66, 0.26, 0.74, 0.18, 0.82, 0.10, 0.90]

    normal_offsets = [
        0.0,
        0.06,
        -0.06,
        0.12,
        -0.12,
        max_normal * 0.45,
        -max_normal * 0.45,
        max_normal * 0.72,
        -max_normal * 0.72,
        max_normal,
        -max_normal,
    ]

    best_x = (sx + dx) / 2.0
    best_y = (sy + dy) / 2.0
    best_score = float("inf")

    for t_pos in t_positions:
        t_used = max(0.08, min(0.92, t_pos + lane_t_bias))
        bx = sx + (dx - sx) * t_used
        by = sy + (dy - sy) * t_used
        for n_off in normal_offsets:
            cx = bx + nx * n_off
            cy = by + ny * n_off

            lx = cx - bw / 2.0
            ly = cy - bh / 2.0
            if lx < 0.02 or ly < 0.02 or lx + bw > slide_w - 0.02 or ly + bh > slide_h - 0.02:
                continue

            score = 0.0
            obstacle_overlaps = overlap_count_with_obstacles(cx, cy, bw, bh, obstacle_boxes)
            label_overlaps = overlap_count_with_labels(cx, cy, bw, bh, placed_labels)
            obstacle_area = overlap_area_with_obstacles(cx, cy, bw, bh, obstacle_boxes)
            label_area = overlap_area_with_labels(cx, cy, bw, bh, placed_labels)

            score += obstacle_overlaps * 14000.0
            score += label_overlaps * 17000.0
            score += obstacle_area * 32000.0
            score += label_area * 38000.0
            score += intersection_penalty(cx, cy, intersections)
            score += abs(n_off) * (80.0 if intersections else 260.0)

            # For crossed edges, pushing away from center is better than occupying the crossing point.
            if intersections:
                score += max(0.0, 0.20 - abs(t_used - 0.5)) * 1800.0
            else:
                score += abs(t_used - 0.5) * 12.0

            dist_start = ((cx - sx) ** 2 + (cy - sy) ** 2) ** 0.5
            dist_end = ((cx - dx) ** 2 + (cy - dy) ** 2) ** 0.5
            endpoint_clearance = min(dist_start, dist_end)
            endpoint_clearance_min = min(0.36, max(0.16, length * 0.42))
            if endpoint_clearance < endpoint_clearance_min:
                score += (endpoint_clearance_min - endpoint_clearance) * 4200.0

            if score < best_score:
                best_score = score
                best_x = cx
                best_y = cy

    return best_x, best_y


def normalize_color_token(value: str | None) -> str | None:
    if not value:
        return None
    text = str(value).strip().strip("\"'")
    matched = re.fullmatch(r"#?([0-9A-Fa-f]{6})", text)
    if matched:
        return matched.group(1).upper()
    named = {
        "black": "000000",
        "white": "FFFFFF",
        "red": "FF0000",
        "green": "008000",
        "blue": "0000FF",
        "yellow": "FFFF00",
        "orange": "FFA500",
        "gray": "808080",
        "grey": "808080",
    }
    return named.get(text.lower())


def parse_style_attr(text: str) -> dict[str, str]:
    out: dict[str, str] = {}
    for part in text.split(";"):
        if ":" not in part:
            continue
        key, value = part.split(":", 1)
        k = key.strip().lower()
        v = value.strip()
        if not k:
            continue
        out[k] = v
    return out


def parse_html_segments(text: str) -> list[tuple[str, dict[str, Any]]]:
    tokens = re.split(r"(<[^>]+>)", text.replace("\r\n", "\n"))
    base_style: dict[str, Any] = {
        "bold": False,
        "italic": False,
        "underline": False,
        "strike": False,
        "code": False,
        "color": None,
    }
    stack: list[tuple[str, dict[str, Any]]] = [("root", base_style)]
    out: list[tuple[str, dict[str, Any]]] = []

    for token in tokens:
        if token == "":
            continue

        if token.startswith("<") and token.endswith(">"):
            if re.match(r"^<\s*br\s*/?\s*>$", token, flags=re.IGNORECASE):
                out.append(("\n", dict(stack[-1][1])))
                continue

            close = re.match(r"^<\s*/\s*([a-zA-Z0-9]+)\s*>$", token)
            if close:
                tag = close.group(1).lower()
                while len(stack) > 1:
                    popped_tag, _ = stack.pop()
                    if popped_tag == tag:
                        break
                continue

            open_m = re.match(r"^<\s*([a-zA-Z0-9]+)([^>]*)>$", token)
            if open_m:
                tag = open_m.group(1).lower()
                attrs = open_m.group(2) or ""
                style = dict(stack[-1][1])

                if tag in {"b", "strong"}:
                    style["bold"] = True
                if tag in {"i", "em"}:
                    style["italic"] = True
                if tag in {"u", "ins"}:
                    style["underline"] = True
                if tag in {"s", "strike", "del"}:
                    style["strike"] = True
                if tag == "code":
                    style["code"] = True

                color_m = re.search(
                    r"""(?:^|\s)color\s*=\s*(?:"([^"]+)"|'([^']+)'|([^\s>]+))""",
                    attrs,
                    flags=re.IGNORECASE,
                )
                if color_m:
                    style["color"] = normalize_color_token(color_m.group(1) or color_m.group(2) or color_m.group(3))

                style_m = re.search(
                    r"""(?:^|\s)style\s*=\s*(?:"([^"]*)"|'([^']*)')""",
                    attrs,
                    flags=re.IGNORECASE,
                )
                if style_m:
                    style_map = parse_style_attr(style_m.group(1) or style_m.group(2) or "")
                    weight = style_map.get("font-weight", "").lower()
                    if "bold" in weight:
                        style["bold"] = True
                    if weight.isdigit() and int(weight) >= 600:
                        style["bold"] = True
                    if "italic" in style_map.get("font-style", "").lower():
                        style["italic"] = True
                    deco = style_map.get("text-decoration", "").lower()
                    if "underline" in deco:
                        style["underline"] = True
                    if "line-through" in deco:
                        style["strike"] = True
                    css_color = normalize_color_token(style_map.get("color"))
                    if css_color:
                        style["color"] = css_color

                stack.append((tag, style))
                continue

            continue

        out.append((token, dict(stack[-1][1])))

    return out


def strip_markdown_wrapper(text: str) -> str:
    stripped = text.strip()
    if len(stripped) >= 2 and stripped.startswith("`") and stripped.endswith("`"):
        return stripped[1:-1]
    return text


def parse_markdown_segments(text: str, base_style: dict[str, Any]) -> list[tuple[str, dict[str, Any]]]:
    if text == "":
        return [("", dict(base_style))]

    if base_style.get("code"):
        return [(text, dict(base_style))]

    normalized = strip_markdown_wrapper(text)
    token_re = re.compile(r"\*\*[\s\S]+?\*\*|__[\s\S]+?__|\*[\s\S]+?\*|_[\s\S]+?_|~~[\s\S]+?~~|`[\s\S]+?`")
    out: list[tuple[str, dict[str, Any]]] = []
    pos = 0
    for matched in token_re.finditer(normalized):
        start, end = matched.span()
        if start > pos:
            out.append((normalized[pos:start], dict(base_style)))
        token = matched.group(0)
        style = dict(base_style)
        inner = token
        if token.startswith("**") and token.endswith("**"):
            style["bold"] = True
            inner = token[2:-2]
        elif token.startswith("__") and token.endswith("__"):
            style["bold"] = True
            inner = token[2:-2]
        elif token.startswith("*") and token.endswith("*"):
            style["italic"] = True
            inner = token[1:-1]
        elif token.startswith("_") and token.endswith("_"):
            style["italic"] = True
            inner = token[1:-1]
        elif token.startswith("~~") and token.endswith("~~"):
            style["strike"] = True
            inner = token[2:-2]
        elif token.startswith("`") and token.endswith("`"):
            style["code"] = True
            inner = token[1:-1]
        out.append((inner, style))
        pos = end

    if pos < len(normalized):
        out.append((normalized[pos:], dict(base_style)))

    return out if out else [(normalized, dict(base_style))]


def build_rich_paragraphs(text: str) -> list[list[tuple[str, dict[str, Any]]]]:
    html_segments = parse_html_segments(text)
    segments: list[tuple[str, dict[str, Any]]] = []
    for seg_text, seg_style in html_segments:
        segments.extend(parse_markdown_segments(seg_text, seg_style))

    paragraphs: list[list[tuple[str, dict[str, Any]]]] = [[]]
    for seg_text, seg_style in segments:
        parts = seg_text.split("\n")
        for i, part in enumerate(parts):
            if part != "":
                paragraphs[-1].append((part, dict(seg_style)))
            if i < len(parts) - 1:
                paragraphs.append([])

    if not paragraphs:
        return [[]]
    return paragraphs


def render_rich_text(
    tf: Any,
    text: str,
    *,
    font_family: str,
    font_size: float,
    color: str,
    bold: bool,
    align: str,
) -> None:
    paragraphs = build_rich_paragraphs(text)
    if not paragraphs:
        paragraphs = [[]]

    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "left":
            p.alignment = PP_ALIGN.LEFT
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.CENTER

        if not runs:
            run = p.add_run()
            run.text = ""
            run.font.name = font_family
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = to_rgb(color)
            continue

        for run_text, run_style in runs:
            if run_text == "":
                continue
            run = p.add_run()
            run.text = run_text
            run.font.name = "Consolas" if run_style.get("code") else font_family
            run.font.size = Pt(font_size)
            run.font.bold = bool(bold or run_style.get("bold"))
            run.font.italic = bool(run_style.get("italic"))
            run.font.underline = bool(run_style.get("underline"))
            run.font.strike = bool(run_style.get("strike"))
            run_color = run_style.get("color") or color
            run.font.color.rgb = to_rgb(str(run_color))


def apply_shape_style(shape: Any, style: dict[str, Any]) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = to_rgb(style.get("fill", "F8FAFC"))

    shape.line.color.rgb = to_rgb(style.get("stroke", "0F172A"))
    shape.line.width = Pt(float(style.get("strokeWidth", 1.0)))


def override_special_node_box(shape_name: str, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    if shape_name == "forkBar":
        bar_w = clampf(w * 0.20, 0.05, 0.10)
        bar_h = clampf(h, 0.36, 0.92)
        cx = x + w / 2.0
        cy = y + h / 2.0
        return cx - bar_w / 2.0, cy - bar_h / 2.0, bar_w, bar_h

    if shape_name == "filledCircle":
        size = clampf(min(w, h), 0.08, 0.18)
        cx = x + w / 2.0
        cy = y + h / 2.0
        return cx - size / 2.0, cy - size / 2.0, size, size

    if shape_name == "smallCircle":
        size = clampf(min(w, h), 0.06, 0.12)
        cx = x + w / 2.0
        cy = y + h / 2.0
        return cx - size / 2.0, cy - size / 2.0, size, size

    if shape_name == "framedCircle":
        size = clampf(min(w, h), 0.09, 0.18)
        cx = x + w / 2.0
        cy = y + h / 2.0
        return cx - size / 2.0, cy - size / 2.0, size, size

    return x, y, w, h


def add_vertical_inner_line(slide: Any, x: float, y: float, w: float, h: float, *, color: str, width_pt: float) -> None:
    line_x = x + clampf(w * 0.16, 0.04, 0.18)
    seg = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(line_x),
        Inches(y + 0.04),
        Inches(line_x),
        Inches(y + h - 0.04),
    )
    seg.line.color.rgb = to_rgb(color)
    seg.line.width = Pt(max(0.6, width_pt))


def add_horizontal_inner_line(slide: Any, x: float, y: float, w: float, *, color: str, width_pt: float) -> None:
    line_y = y + 0.10
    seg = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x + 0.04),
        Inches(line_y),
        Inches(x + w - 0.04),
        Inches(line_y),
    )
    seg.line.color.rgb = to_rgb(color)
    seg.line.width = Pt(max(0.6, width_pt))


def decorate_special_node(
    slide: Any,
    shape: Any,
    shape_name: str,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style: dict[str, Any],
) -> None:
    stroke = style.get("stroke", "0F172A")
    stroke_w = float(style.get("strokeWidth", 1.0))

    if shape_name == "filledCircle":
        shape.fill.solid()
        shape.fill.fore_color.rgb = to_rgb(stroke)
        shape.line.color.rgb = to_rgb(stroke)
        shape.line.width = Pt(max(0.8, stroke_w))
        return

    if shape_name == "framedCircle":
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + w * 0.32),
            Inches(y + h * 0.32),
            Inches(max(0.03, w * 0.36)),
            Inches(max(0.03, h * 0.36)),
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = to_rgb(stroke)
        inner.line.fill.background()
        return

    if shape_name == "linedRect":
        add_vertical_inner_line(slide, x, y, w, h, color=stroke, width_pt=stroke_w)
        return

    if shape_name == "linedDocument":
        add_vertical_inner_line(slide, x, y, w, h, color=stroke, width_pt=stroke_w)
        return

    if shape_name == "windowPane":
        add_vertical_inner_line(slide, x, y, w, h, color=stroke, width_pt=stroke_w)
        add_horizontal_inner_line(slide, x, y, w, color=stroke, width_pt=stroke_w)
        return

    if shape_name == "framedRect":
        add_vertical_inner_line(slide, x, y, w, h, color=stroke, width_pt=stroke_w)
        right_line_x = x + w - clampf(w * 0.16, 0.04, 0.18)
        seg = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(right_line_x),
            Inches(y + 0.04),
            Inches(right_line_x),
            Inches(y + h - 0.04),
        )
        seg.line.color.rgb = to_rgb(stroke)
        seg.line.width = Pt(max(0.6, stroke_w))
        return

    if shape_name == "stackedRect":
        # draw trailing stacked layers
        for offset in (0.05, 0.10):
            back = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + offset),
                Inches(y - offset),
                Inches(max(0.05, w)),
                Inches(max(0.05, h)),
            )
            back.fill.solid()
            back.fill.fore_color.rgb = to_rgb(style.get("fill", "F8FAFC"))
            back.line.color.rgb = to_rgb(stroke)
            back.line.width = Pt(max(0.6, stroke_w))
        return

def set_shape_text(
    shape: Any,
    text: str,
    *,
    font_family: str,
    font_size: float,
    color: str,
    bold: bool = False,
    align: str = "center",
    vertical: str = "middle",
    auto_fit: bool = True,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_fit else MSO_AUTO_SIZE.NONE

    if vertical == "top":
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)

    render_rich_text(
        tf,
        text,
        font_family=font_family,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
    )


def render(
    ir: dict[str, Any],
    output: Path,
    patch_text: str | None,
    slide_size: str | None,
    edge_routing: str | None = "straight",
    append_to: Path | None = None,
) -> None:
    append_target = append_to.resolve() if append_to else None
    if append_target and append_target.exists():
        prs = Presentation(str(append_target))
        base_w = emu_to_in(float(prs.slide_width))
        base_h = emu_to_in(float(prs.slide_height))
        max_w = base_w
        max_h = base_h
    else:
        prs = Presentation()
        base_w, base_h, max_w, max_h = parse_slide_size(slide_size)

    slide_w, slide_h, scale, offset_x, offset_y = resolve_slide_size_and_scale(
        ir,
        base_slide_w=base_w,
        base_slide_h=base_h,
        max_slide_w=max_w,
        max_slide_h=max_h,
    )
    if not (append_target and append_target.exists()):
        prs.slide_width = Inches(slide_w)
        prs.slide_height = Inches(slide_h)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    config = ir.get("config", {})
    font_family = config.get("fontFamily", "Yu Gothic UI")
    routing_mode = (edge_routing or "straight").strip().lower()
    if routing_mode not in {"straight", "elbow"}:
        routing_mode = "straight"

    node_map: dict[str, dict[str, Any]] = {node["id"]: node for node in ir.get("nodes", [])}
    subgraph_map: dict[str, dict[str, Any]] = {
        subgraph["id"]: subgraph for subgraph in ir.get("subgraphs", []) if isinstance(subgraph.get("id"), str)
    }
    node_shape_map: dict[str, Any] = {}
    node_icon_map: dict[str, Path] = {}
    subgraph_shape_map: dict[str, Any] = {}
    node_box_map: dict[str, tuple[float, float, float, float]] = {}
    obstacle_boxes: list[tuple[float, float, float, float]] = []

    subgraph_title_items: list[dict[str, Any]] = []

    for subgraph in ir.get("subgraphs", []):
        x, y = transform(float(subgraph["x"]), float(subgraph["y"]), scale, offset_x, offset_y)
        w = px_to_in(float(subgraph["width"])) * scale
        h = px_to_in(float(subgraph["height"])) * scale

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(max(w, 0.05)),
            Inches(max(h, 0.05)),
        )

        style = subgraph.get("style", {})
        shape.fill.solid()
        shape.fill.fore_color.rgb = to_rgb(style.get("fill", "F8FAFC"))
        shape.fill.transparency = 0.9
        shape.line.color.rgb = to_rgb(style.get("stroke", "64748B"))
        shape.line.width = Pt(float(style.get("strokeWidth", 1.0)))
        if style.get("dash") == "dash":
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        subgraph_id = str(subgraph.get("id", "")).strip()
        if subgraph_id:
            subgraph_shape_map[subgraph_id] = shape

        title_text = str(subgraph.get("title", "")).strip()
        if title_text:
            title_font = max(7.0, scaled_font_size(9.0, scale))
            title_x = x + 0.06
            title_w = max(0.40, w - 0.12)
            wrapped_lines = estimate_wrapped_line_count(title_text, title_w, title_font)
            desired_h = max(0.36, (title_font * 1.30 * wrapped_lines) / 72.0 + 0.12)

            member_nodes = [
                node_map.get(node_id)
                for node_id in subgraph.get("nodeIds", [])
                if isinstance(node_id, str) and node_map.get(node_id) is not None
            ]
            if member_nodes:
                min_member_y = min(
                    transform(float(node["x"]), float(node["y"]), scale, offset_x, offset_y)[1] for node in member_nodes
                )
                top_limit = y + 0.02
                bottom_limit = min(min_member_y - 0.08, y + h - 0.04)
                if bottom_limit <= top_limit + 0.12:
                    bottom_limit = min(y + h - 0.04, top_limit + 0.24)

                title_y = top_limit
                available_h = max(0.20, bottom_limit - top_limit)
                title_h = available_h
            else:
                title_y = y + 0.02
                available_h = max(0.24, min(y + h - 0.04 - title_y, 1.40))
                title_h = available_h

            title_h = min(title_h, max(0.20, y + h - 0.04 - title_y))
            wrapped_lines = estimate_wrapped_line_count(title_text, title_w, title_font)
            max_font_by_height = (max(0.16, title_h) * 72.0 - 0.10) / max(1.28 * max(1, wrapped_lines), 1e-6)
            title_font = min(title_font, max(5.5, max_font_by_height))
            title_text = compact_text_to_fit_box(title_text, title_w, title_h, title_font)
            subgraph_title_items.append(
                {
                    "text": title_text,
                    "x": title_x,
                    "y": title_y,
                    "w": title_w,
                    "h": max(0.20, title_h),
                    "font": title_font,
                    "fill": style.get("fill", "F8FAFC"),
                }
            )
            obstacle_boxes.append((title_x, title_y, title_w, max(0.20, title_h)))

    for node in ir.get("nodes", []):
        node_id = str(node.get("id", "")).strip()
        icon_raw = str(node.get("icon", "")).strip()
        if not node_id or not icon_raw:
            continue
        icon_path = resolve_icon_png(icon_raw, str(node.get("label", "")), node_id, size_px=128)
        if icon_path is not None:
            node_icon_map[node_id] = icon_path

    for node in ir.get("nodes", []):
        x, y = transform(float(node["x"]), float(node["y"]), scale, offset_x, offset_y)
        w = px_to_in(float(node["width"])) * scale
        h = px_to_in(float(node["height"])) * scale
        shape_name = str(node.get("shape", "rect"))
        is_junction = bool(node.get("isJunction", False))
        style = dict(node.get("style", {}))
        theme = infer_arch_theme(node)
        if theme:
            if str(style.get("fill", "F8FAFC")).strip().upper() == "F8FAFC":
                style["fill"] = theme["fill"]
            if str(style.get("stroke", "0F172A")).strip().upper() == "0F172A":
                style["stroke"] = theme["stroke"]
            if str(style.get("text", "0F172A")).strip().upper() == "0F172A":
                style["text"] = theme["text"]

        if is_junction:
            # Junction is a logical connection point. Keep it invisible in pptx.
            anchor = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(max(w, 0.02)),
                Inches(max(h, 0.02)),
            )
            anchor.fill.background()
            anchor.line.fill.background()
            node_shape_map[node["id"]] = anchor
            node_box_map[node["id"]] = (x, y, w, h)
            continue

        display_label = str(node.get("label", node.get("id", "")))
        if shape_name in {"forkBar", "filledCircle", "smallCircle", "framedCircle"}:
            display_label = ""

        node_id = str(node.get("id", "")).strip()
        icon_path = node_icon_map.get(node_id)
        if icon_path is not None and shape_name not in {"forkBar", "filledCircle", "smallCircle", "framedCircle"}:
            # Render icon nodes as one grouped object:
            # transparent anchor (for connectors) + large icon + label below.
            node_group = slide.shapes.add_group_shape()
            anchor = node_group.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(max(w, 0.05)),
                Inches(max(h, 0.05)),
            )
            anchor.fill.background()
            anchor.line.fill.background()

            label_band_h = clampf(max(0.30, h * 0.30), 0.28, max(0.40, h * 0.48))
            icon_area_h = max(0.10, h - label_band_h - 0.05)
            icon_size = clampf(min(w * 0.82, icon_area_h * 0.90), 0.22, min(1.30, max(0.22, w - 0.08)))
            ix = x + (w - icon_size) / 2.0
            iy = y + max(0.02, (icon_area_h - icon_size) / 2.0 + 0.02)

            try:
                node_group.shapes.add_picture(str(icon_path), Inches(ix), Inches(iy), Inches(icon_size), Inches(icon_size))
            except Exception:
                # Fallback to default node rendering when icon drawing fails.
                pass

            if display_label:
                tx = x + 0.03
                tw = max(0.10, w - 0.06)
                th = max(0.20, label_band_h - 0.04)
                ty = y + h - th - 0.02
                label_box = node_group.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
                label_box.fill.background()
                label_box.line.fill.background()
                set_shape_text(
                    label_box,
                    display_label,
                    font_family=font_family,
                    font_size=scaled_font_size(float(style.get("fontSize", 14)), scale),
                    color=style.get("text", "0F172A"),
                    bold=bool(style.get("bold", False)),
                    align="center",
                    vertical="top",
                    auto_fit=True,
                )

            node_shape_map[node["id"]] = anchor
            node_box_map[node["id"]] = (x, y, w, h)
            obstacle_boxes.append((x, y, w, h))
            continue

        is_stacked_rect = shape_name == "stackedRect"

        if is_stacked_rect:
            decorate_special_node(slide, None, "stackedRect", x=x, y=y, w=w, h=h, style=style)

        x, y, w, h = override_special_node_box(shape_name, x, y, w, h)

        shape = slide.shapes.add_shape(
            SHAPE_MAP.get(shape_name, MSO_SHAPE.RECTANGLE),
            Inches(x),
            Inches(y),
            Inches(max(w, 0.05)),
            Inches(max(h, 0.05)),
        )

        apply_shape_style(shape, style)
        if not is_stacked_rect:
            decorate_special_node(slide, shape, shape_name, x=x, y=y, w=w, h=h, style=style)

        set_shape_text(
            shape,
            display_label,
            font_family=font_family,
            font_size=scaled_font_size(float(style.get("fontSize", 14)), scale),
            color=style.get("text", "0F172A"),
            bold=bool(style.get("bold", False)),
            align="center",
            vertical="middle",
        )

        node_shape_map[node["id"]] = shape
        node_box_map[node["id"]] = (x, y, w, h)
        obstacle_boxes.append((x, y, w, h))

    # Draw subgraph titles after nodes so they remain visible and are not hidden behind objects.
    for title_item in subgraph_title_items:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(title_item["x"]),
            Inches(title_item["y"]),
            Inches(title_item["w"]),
            Inches(title_item["h"]),
        )
        box.fill.background()
        box.line.fill.background()
        set_shape_text(
            box,
            title_item["text"],
            font_family=font_family,
            font_size=title_item["font"],
            color="000000",
            bold=True,
            align="left",
            vertical="top",
            auto_fit=True,
        )

    edge_items: list[dict[str, Any]] = []
    pending_edges: list[dict[str, Any]] = []

    for edge in ir.get("edges", []):
        src = node_map.get(edge.get("from"))
        dst = node_map.get(edge.get("to"))
        src_shape = node_shape_map.get(edge.get("from"))
        dst_shape = node_shape_map.get(edge.get("to"))
        if not src or not dst:
            continue

        style = edge.get("style", {})
        src_id = str(edge.get("from", ""))
        dst_id = str(edge.get("to", ""))

        src_anchor = src
        dst_anchor = dst
        src_connect_shape = src_shape
        dst_connect_shape = dst_shape
        src_anchor_id = src_id
        dst_anchor_id = dst_id

        if bool(style.get("startViaGroup")):
            src_group_id = str(src.get("subgraphId", "")).strip()
            group = subgraph_map.get(src_group_id)
            if group is not None:
                src_anchor = group
                src_connect_shape = subgraph_shape_map.get(src_group_id)
                src_anchor_id = src_group_id

        if bool(style.get("endViaGroup")):
            dst_group_id = str(dst.get("subgraphId", "")).strip()
            group = subgraph_map.get(dst_group_id)
            if group is not None:
                dst_anchor = group
                dst_connect_shape = subgraph_shape_map.get(dst_group_id)
                dst_anchor_id = dst_group_id

        is_self_loop = src_id == dst_id and not bool(style.get("startViaGroup")) and not bool(style.get("endViaGroup"))

        if is_self_loop:
            node_box = node_box_map.get(src_id)
            if node_box is None:
                continue
            src_side, dst_side, loop_points, loop_label_anchor = choose_self_loop_geometry(
                node_box, slide_w=slide_w, slide_h=slide_h
            )
        else:
            requested_src_side = side_from_token(style.get("startSide"))
            requested_dst_side = side_from_token(style.get("endSide"))
            strict_hints = bool(src_anchor.get("isJunction", False) or dst_anchor.get("isJunction", False))
            src_side, dst_side = choose_connection_sides_with_hints(
                src_anchor,
                dst_anchor,
                hinted_src_side=requested_src_side,
                hinted_dst_side=requested_dst_side,
                strict_hints=strict_hints,
            )
            loop_points = None
            loop_label_anchor = None

        pending_edges.append(
            {
                "edge": edge,
                "srcAnchor": src_anchor,
                "dstAnchor": dst_anchor,
                "srcShape": src_connect_shape,
                "dstShape": dst_connect_shape,
                "srcId": src_id,
                "dstId": dst_id,
                "srcAnchorId": src_anchor_id,
                "dstAnchorId": dst_anchor_id,
                "isSelfLoop": is_self_loop,
                "srcSide": src_side,
                "dstSide": dst_side,
                "loopPoints": loop_points,
                "loopLabelAnchor": loop_label_anchor,
                "srcOffsetPx": 0.0,
                "dstOffsetPx": 0.0,
            }
        )

    side_groups: dict[tuple[str, int], list[tuple[int, bool]]] = {}
    for idx, item in enumerate(pending_edges):
        if item["isSelfLoop"]:
            continue
        side_groups.setdefault((item["srcAnchorId"], int(item["srcSide"])), []).append((idx, True))
        side_groups.setdefault((item["dstAnchorId"], int(item["dstSide"])), []).append((idx, False))

    for (node_id, side), members in side_groups.items():
        if len(members) <= 1:
            continue
        node = node_map.get(node_id) or subgraph_map.get(node_id)
        if not node:
            continue
        if bool(node.get("isJunction", False)):
            continue
        span = float(node["width"]) if side in {TOP, BOTTOM} else float(node["height"])
        sorted_members = sorted(
            members,
            key=lambda item: side_sort_axis(
                side,
                pending_edges[item[0]]["srcAnchor"],
                pending_edges[item[0]]["dstAnchor"],
                source_side=item[1],
            ),
        )
        offsets = lane_offsets(len(sorted_members), span)
        for offset_idx, (edge_idx, is_source) in enumerate(sorted_members):
            if is_source:
                pending_edges[edge_idx]["srcOffsetPx"] = offsets[offset_idx]
            else:
                pending_edges[edge_idx]["dstOffsetPx"] = offsets[offset_idx]

    for pending in pending_edges:
        edge = pending["edge"]
        src = pending["srcAnchor"]
        dst = pending["dstAnchor"]
        src_shape = pending["srcShape"]
        dst_shape = pending["dstShape"]
        is_self_loop = bool(pending["isSelfLoop"])
        src_side = int(pending["srcSide"])
        dst_side = int(pending["dstSide"])
        loop_points = pending["loopPoints"]
        loop_label_anchor = pending["loopLabelAnchor"]
        src_offset_px = float(pending["srcOffsetPx"])
        dst_offset_px = float(pending["dstOffsetPx"])

        edge_group = slide.shapes.add_group_shape()
        style = edge.get("style", {})
        connector = None
        if is_self_loop and loop_points and len(loop_points) >= 2:
            segments: list[Any] = []
            for idx in range(len(loop_points) - 1):
                p0 = loop_points[idx]
                p1 = loop_points[idx + 1]
                segment = edge_group.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(p0[0]),
                    Inches(p0[1]),
                    Inches(p1[0]),
                    Inches(p1[1]),
                )
                apply_line_style(segment.line, style)
                segments.append(segment)

            if segments:
                start_style = dict(style)
                start_style["endMarker"] = "none"
                start_style["arrow"] = "start"
                end_style = dict(style)
                end_style["startMarker"] = "none"
                end_style["arrow"] = "end"
                set_edge_markers(segments[0], start_style)
                set_edge_markers(segments[-1], end_style)
                connector = segments[-1]
                sx, sy = loop_points[0]
                dx, dy = loop_points[-1]
            else:
                continue
        else:
            src_x, src_y = side_anchor_point(src, src_side, src_offset_px)
            dst_x, dst_y = side_anchor_point(dst, dst_side, dst_offset_px)
            sx, sy = transform(src_x, src_y, scale, offset_x, offset_y)
            dx, dy = transform(dst_x, dst_y, scale, offset_x, offset_y)
            start_marker_name = str(style.get("startMarker", "")).strip()
            end_marker_name = str(style.get("endMarker", "")).strip()
            marker_needs_manual_anchor = start_marker_name in {"circle", "openDiamond"} or end_marker_name in {
                "circle",
                "openDiamond",
            }
            use_manual_anchors = abs(src_offset_px) > 0.8 or abs(dst_offset_px) > 0.8 or marker_needs_manual_anchor

            connector_type = MSO_CONNECTOR.ELBOW if routing_mode == "elbow" else MSO_CONNECTOR.STRAIGHT
            connector = edge_group.shapes.add_connector(
                connector_type,
                Inches(sx),
                Inches(sy),
                Inches(dx),
                Inches(dy),
            )

            if not use_manual_anchors and src_shape is not None and dst_shape is not None:
                try:
                    connector.begin_connect(src_shape, src_side)
                    connector.end_connect(dst_shape, dst_side)
                except Exception:
                    pass

            apply_line_style(connector.line, style)
            set_edge_markers(connector, style)
            if not use_manual_anchors:
                sx = emu_to_in(float(connector.begin_x))
                sy = emu_to_in(float(connector.begin_y))
                dx = emu_to_in(float(connector.end_x))
                dy = emu_to_in(float(connector.end_y))

        start_marker_name = str(style.get("startMarker", "")).strip()
        end_marker_name = str(style.get("endMarker", "")).strip()
        if start_marker_name in {"openDiamond", "circle"}:
            add_endpoint_symbol(
                edge_group,
                start_marker_name,
                at_start=True,
                sx=sx,
                sy=sy,
                dx=dx,
                dy=dy,
                color=str(style.get("color", "1E293B")),
            )
        if end_marker_name in {"openDiamond", "circle"}:
            add_endpoint_symbol(
                edge_group,
                end_marker_name,
                at_start=False,
                sx=sx,
                sy=sy,
                dx=dx,
                dy=dy,
                color=str(style.get("color", "1E293B")),
            )

        edge_items.append(
            {
                "group": edge_group,
                "connector": connector,
                "style": style,
                "label": edge.get("label"),
                "startLabel": edge.get("startLabel"),
                "endLabel": edge.get("endLabel"),
                "sx": sx,
                "sy": sy,
                "dx": dx,
                "dy": dy,
                "intersections": [],
                "selfLoop": is_self_loop,
                "loopLabelAnchor": loop_label_anchor,
            }
        )

    for i in range(len(edge_items)):
        for j in range(i + 1, len(edge_items)):
            i_item = edge_items[i]
            j_item = edge_items[j]
            if i_item.get("selfLoop") or j_item.get("selfLoop"):
                continue
            p = segment_intersection(
                (i_item["sx"], i_item["sy"]),
                (i_item["dx"], i_item["dy"]),
                (j_item["sx"], j_item["sy"]),
                (j_item["dx"], j_item["dy"]),
            )
            if p is None:
                continue
            i_item["intersections"].append(p)
            j_item["intersections"].append(p)

    placed_label_boxes: list[tuple[float, float, float, float]] = []
    labeled_items = [item for item in edge_items if item["label"]]
    labeled_items.sort(key=lambda item: len(item["intersections"]), reverse=True)

    segment_totals: dict[tuple[float, float, float, float], int] = {}
    for item in labeled_items:
        key = (
            round(min(item["sx"], item["dx"]), 2),
            round(min(item["sy"], item["dy"]), 2),
            round(max(item["sx"], item["dx"]), 2),
            round(max(item["sy"], item["dy"]), 2),
        )
        segment_totals[key] = segment_totals.get(key, 0) + 1

    segment_seen: dict[tuple[float, float, float, float], int] = {}

    for item in labeled_items:
        label_text = str(item["label"])
        label_font = scaled_font_size(float(item["style"].get("fontSize", 11)), scale)
        bw, bh = estimate_label_box(label_text, label_font)
        seg_key = (
            round(min(item["sx"], item["dx"]), 2),
            round(min(item["sy"], item["dy"]), 2),
            round(max(item["sx"], item["dx"]), 2),
            round(max(item["sy"], item["dy"]), 2),
        )
        lane_index = segment_seen.get(seg_key, 0)
        segment_seen[seg_key] = lane_index + 1
        lane_total = segment_totals.get(seg_key, 1)
        lane_center = (lane_total - 1) / 2.0
        lane_t_bias = (lane_index - lane_center) * 0.08

        if item.get("selfLoop") and item.get("loopLabelAnchor") is not None:
            base_x, base_y = item["loopLabelAnchor"]
            offsets = [(0.0, 0.0), (0.0, -0.10), (0.10, 0.0), (-0.10, 0.0), (0.0, 0.10)]
            lx = base_x
            ly = base_y
            best_penalty = float("inf")
            for ox, oy in offsets:
                cx = clampf(base_x + ox, 0.06 + bw / 2.0, slide_w - 0.06 - bw / 2.0)
                cy = clampf(base_y + oy, 0.06 + bh / 2.0, slide_h - 0.06 - bh / 2.0)
                penalty = 0.0
                penalty += overlap_count_with_obstacles(cx, cy, bw, bh, obstacle_boxes) * 9000.0
                penalty += overlap_count_with_labels(cx, cy, bw, bh, placed_label_boxes) * 14000.0
                if penalty < best_penalty:
                    best_penalty = penalty
                    lx = cx
                    ly = cy
        else:
            lx, ly = choose_label_center(
                item["sx"],
                item["sy"],
                item["dx"],
                item["dy"],
                bw,
                bh,
                intersections=item["intersections"],
                obstacle_boxes=obstacle_boxes,
                placed_labels=placed_label_boxes,
                slide_w=slide_w,
                slide_h=slide_h,
                lane_t_bias=lane_t_bias,
            )

        box = item["group"].shapes.add_textbox(Inches(lx - bw / 2.0), Inches(ly - bh / 2.0), Inches(bw), Inches(bh))
        box.fill.background()
        box.line.fill.background()

        tf = box.text_frame
        tf.clear()
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)

        render_rich_text(
            tf,
            label_text,
            font_family=font_family,
            font_size=label_font,
            color=item["style"].get("color", "1E293B"),
            bold=False,
            align="center",
        )

        placed_label_boxes.append((lx - bw / 2.0, ly - bh / 2.0, bw, bh))

    for item in edge_items:
        for side in ("start", "end"):
            endpoint_text_raw = item["startLabel"] if side == "start" else item["endLabel"]
            endpoint_text = str(endpoint_text_raw).strip() if endpoint_text_raw is not None else ""
            if endpoint_text == "":
                continue

            sx = float(item["sx"])
            sy = float(item["sy"])
            dx = float(item["dx"])
            dy = float(item["dy"])
            vx = dx - sx
            vy = dy - sy
            length = (vx * vx + vy * vy) ** 0.5
            if length < 1e-6:
                continue

            tx = vx / length
            ty = vy / length
            nx = -ty
            ny = tx
            t = 0.14 if side == "start" else 0.86
            base_x = sx + vx * t
            base_y = sy + vy * t
            offset = 0.11
            cx = base_x + nx * offset
            cy = base_y + ny * offset

            label_font = scaled_font_size(float(item["style"].get("fontSize", 10)), scale)
            bw, bh = estimate_label_box(endpoint_text, max(7.0, label_font - 0.8))
            cx = clampf(cx, 0.04 + bw / 2.0, slide_w - 0.04 - bw / 2.0)
            cy = clampf(cy, 0.04 + bh / 2.0, slide_h - 0.04 - bh / 2.0)

            box = item["group"].shapes.add_textbox(Inches(cx - bw / 2.0), Inches(cy - bh / 2.0), Inches(bw), Inches(bh))
            box.fill.background()
            box.line.fill.background()

            tf = box.text_frame
            tf.clear()
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.margin_left = Inches(0.01)
            tf.margin_right = Inches(0.01)
            tf.margin_top = Inches(0.0)
            tf.margin_bottom = Inches(0.0)

            render_rich_text(
                tf,
                endpoint_text,
                font_family=font_family,
                font_size=max(7.0, label_font - 0.8),
                color=item["style"].get("color", "1E293B"),
                bold=False,
                align="center",
            )
            placed_label_boxes.append((cx - bw / 2.0, cy - bh / 2.0, bw, bh))

    title_text = str(ir.get("meta", {}).get("title", "")).strip()
    if title_text:
        title_shape = slide.shapes.add_textbox(
            Inches(0.20),
            Inches(0.08),
            Inches(max(0.80, slide_w - 0.40)),
            Inches(0.24),
        )
        title_shape.fill.background()
        title_shape.line.fill.background()

        title_frame = title_shape.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        title_frame.margin_left = Inches(0.0)
        title_frame.margin_right = Inches(0.0)
        title_frame.margin_top = Inches(0.0)
        title_frame.margin_bottom = Inches(0.0)

        render_rich_text(
            title_frame,
            title_text,
            font_family=font_family,
            font_size=10.0,
            color="334155",
            bold=True,
            align="left",
        )

    embed = {
        "sourceMmd": ir.get("meta", {}).get("source", ""),
        "patchYaml": patch_text,
        "version": 1,
    }
    payload = "\n".join(
        [
            "__MMD2PPTX_EMBED_START__",
            json.dumps(embed, ensure_ascii=False, indent=2),
            "__MMD2PPTX_EMBED_END__",
        ]
    )

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = payload

    save_path = append_target if append_target else output
    save_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(save_path))


def main() -> int:
    parser = argparse.ArgumentParser(description="Render mmd2pptx IR to PowerPoint using python-pptx")
    parser.add_argument("--ir", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--patch", type=Path, default=None)
    parser.add_argument("--slide-size", type=str, default="16:9")
    parser.add_argument("--edge-routing", type=str, default="straight", choices=["straight", "elbow"])
    parser.add_argument("--append-to", type=Path, default=None)
    args = parser.parse_args()

    ir = json.loads(args.ir.read_text(encoding="utf-8"))
    patch_text = args.patch.read_text(encoding="utf-8") if args.patch else None
    render(ir, args.output, patch_text, args.slide_size, args.edge_routing, args.append_to)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
