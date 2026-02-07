#!/usr/bin/env python3
from __future__ import annotations

import argparse
import html
import json
import math
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

PX_PER_INCH = 96.0
EMU_PER_INCH = 914400.0

PARTICIPANT_KINDS = {
    "participant",
    "actor",
    "boundary",
    "control",
    "entity",
    "database",
    "collections",
    "queue",
}

ARROW_PATTERNS = [
    "<<-->>",
    "<<->>",
    "--|\\",
    "--|/",
    "--\\\\",
    "--//",
    "/|--",
    "\\|--",
    "//--",
    "\\\\--",
    "-->>",
    "->>",
    "-|\\",
    "-|/",
    "-\\\\",
    "-//",
    "/|-",
    "\\|-",
    "//-",
    "\\\\-",
    "--x",
    "-x",
    "--)",
    "-)",
    "-->",
    "->",
]
ARROW_PATTERNS.sort(key=len, reverse=True)

DOTTED_ARROWS = {
    "-->",
    "-->>",
    "<<-->>",
    "--x",
    "--)",
    "--|\\",
    "--|/",
    "--\\\\",
    "--//",
    "/|--",
    "\\|--",
    "//--",
    "\\\\--",
}

HALF_ARROW_FORWARD = {"-|\\", "-|/", "-\\\\", "-//", "--|\\", "--|/", "--\\\\", "--//"}
HALF_ARROW_REVERSE = {"/|-", "\\|-", "//-", "\\\\-", "/|--", "\\|--", "//--", "\\\\--"}
ENTITY_TOKEN_RE = re.compile(r"#(x[0-9A-Fa-f]+|[0-9]+|[A-Za-z][A-Za-z0-9]+);")


def clampf(value: float, min_v: float, max_v: float) -> float:
    return max(min_v, min(max_v, value))


def to_rgb(value: str) -> RGBColor:
    text = (value or "000000").strip().lstrip("#")
    if len(text) != 6:
        text = "000000"
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def parse_css_color(value: str, fallback: str = "E2E8F0") -> str:
    text = (value or "").strip()
    if not text:
        return fallback

    hex_match = re.fullmatch(r"#?([0-9a-fA-F]{6})", text)
    if hex_match:
        return hex_match.group(1).upper()

    rgb_match = re.fullmatch(r"rgba?\(([^)]+)\)", text, flags=re.IGNORECASE)
    if rgb_match:
        parts = [p.strip() for p in rgb_match.group(1).split(",")]
        if len(parts) >= 3:
            try:
                r = int(float(parts[0]))
                g = int(float(parts[1]))
                b = int(float(parts[2]))
                r = max(0, min(255, r))
                g = max(0, min(255, g))
                b = max(0, min(255, b))
                return f"{r:02X}{g:02X}{b:02X}"
            except Exception:
                pass

    named = {
        "transparent": "FFFFFF",
        "white": "FFFFFF",
        "black": "000000",
        "gray": "CBD5E1",
        "grey": "CBD5E1",
        "red": "FCA5A5",
        "green": "86EFAC",
        "blue": "93C5FD",
        "yellow": "FDE68A",
        "orange": "FDBA74",
        "purple": "C4B5FD",
        "aqua": "67E8F9",
        "teal": "5EEAD4",
    }
    return named.get(text.lower(), fallback)


def set_line_markers(connector: Any, start: str = "none", end: str = "none") -> None:
    ln = connector._element.spPr.get_or_add_ln()
    for child in list(ln):
        if child.tag in {qn("a:headEnd"), qn("a:tailEnd")}:
            ln.remove(child)

    if start != "none":
        head = OxmlElement("a:headEnd")
        head.set("type", start)
        ln.append(head)

    if end != "none":
        tail = OxmlElement("a:tailEnd")
        tail.set("type", end)
        ln.append(tail)


def text_units(text: str) -> float:
    units = 0.0
    for ch in text:
        if ch.isspace():
            units += 0.35
            continue
        code = ord(ch)
        if code >= 0x3000:
            units += 1.75
            continue
        units += 1.0
    return units


def split_lines(text: str) -> list[str]:
    return [line for line in text.replace("\r\n", "\n").split("\n") if line.strip()]


def clean_text(raw: str) -> str:
    out = raw.strip()
    quoted = re.fullmatch(r"(['\"])(.*)\1", out)
    if quoted:
        out = quoted.group(2)
    out = out.replace("<br/>", "\n").replace("<br />", "\n").replace("<br>", "\n")
    out = decode_entity_tokens(out)
    out = html.unescape(out)
    return out.strip()


def normalize_participant_id(raw: str) -> str:
    token = clean_text(raw)
    if not token:
        token = "participant"
    token = re.sub(r"\s+", "_", token)
    token = re.sub(r"[^A-Za-z0-9_.:-]", "_", token)
    return token


def decode_entity_tokens(text: str) -> str:
    def repl(match: re.Match[str]) -> str:
        token = match.group(1)
        if token.lower().startswith("x"):
            try:
                return chr(int(token[1:], 16))
            except Exception:
                return match.group(0)
        if token.isdigit():
            try:
                return chr(int(token, 10))
            except Exception:
                return match.group(0)
        return html.unescape(f"&{token};")

    return ENTITY_TOKEN_RE.sub(repl, text)


def parse_json_object(raw: str) -> dict[str, Any]:
    candidates = [raw]
    stripped = raw.strip()
    if stripped and not stripped.startswith("{"):
        candidates.append(f"{{{stripped}}}")
    for candidate in candidates:
        try:
            parsed = json.loads(candidate)
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            continue
    return {}


def split_alias(raw: str) -> tuple[str, str | None]:
    matched = re.match(r"^(.*?)\s+as\s+(.+)$", raw, flags=re.IGNORECASE)
    if not matched:
        return raw.strip(), None
    return matched.group(1).strip(), matched.group(2).strip()


def parse_participant_spec(raw: str, *, default_kind: str = "participant") -> tuple[str, str, str] | None:
    text = raw.strip()
    if not text:
        return None

    explicit_kind = default_kind
    kind_match = re.match(r"^(participant|actor|boundary|control|entity|database|collections|queue)\s+(.+)$", text, flags=re.IGNORECASE)
    if kind_match:
        explicit_kind = kind_match.group(1).lower()
        text = kind_match.group(2).strip()

    base, external_alias = split_alias(text)
    config_obj: dict[str, Any] = {}
    config_match = re.match(r"^(.+?)@\{(.*)\}$", base)
    if config_match:
        base = config_match.group(1).strip()
        config_obj = parse_json_object(config_match.group(2).strip())

    if not base and external_alias:
        base = external_alias
        external_alias = None
    if not base:
        return None

    base_label = clean_text(base)
    cfg_alias = clean_text(str(config_obj.get("alias") or ""))
    config_type = str(config_obj.get("type") or "").strip().lower()

    if base.startswith(("'", '"')) and external_alias:
        pid = normalize_participant_id(external_alias)
        label = base_label
    else:
        pid = normalize_participant_id(base)
        label = clean_text(external_alias or cfg_alias or base_label)

    if not pid:
        return None
    if not label:
        label = pid

    kind = explicit_kind.lower() if explicit_kind else "participant"
    if config_type in PARTICIPANT_KINDS:
        kind = config_type
    if kind not in PARTICIPANT_KINDS:
        kind = "participant"

    return kind, pid, label


def parse_participant_decl(line: str) -> tuple[str, str, str] | None:
    if not re.match(r"^(participant|actor|boundary|control|entity|database|collections|queue)\b", line, flags=re.IGNORECASE):
        return None
    return parse_participant_spec(line)


def parse_create_decl(line: str) -> tuple[str, str, str] | None:
    matched = re.match(r"^create\s+(.+)$", line, flags=re.IGNORECASE)
    if not matched:
        return None
    return parse_participant_spec(matched.group(1), default_kind="participant")


def parse_destroy_decl(line: str) -> tuple[str, str, str] | None:
    matched = re.match(r"^destroy\s+(.+)$", line, flags=re.IGNORECASE)
    if not matched:
        return None
    return parse_participant_spec(matched.group(1), default_kind="participant")


def parse_actor_ref(raw: str) -> str:
    spec = parse_participant_spec(raw, default_kind="participant")
    if spec:
        return spec[1]
    return normalize_participant_id(raw)


def split_message_line(line: str) -> tuple[str, str]:
    if ":" not in line:
        return line.strip(), ""
    left, right = line.split(":", 1)
    return left.strip(), clean_text(right)


def parse_message(line: str) -> dict[str, Any] | None:
    head, label = split_message_line(line)
    if not head:
        return None

    for arrow in ARROW_PATTERNS:
        idx = head.find(arrow)
        if idx < 0:
            continue

        left = head[:idx].strip()
        right = head[idx + len(arrow) :].strip()

        central_from = False
        central_to = False
        if left.endswith("()"):
            central_from = True
            left = left[:-2].strip()
        if right.startswith("()"):
            central_to = True
            right = right[2:].strip()

        suffix = ""
        if right and right[0] in "+-":
            suffix = right[0]
            right = right[1:].strip()

        src = parse_actor_ref(left)
        dst = parse_actor_ref(right)
        if not src or not dst:
            return None

        return {
            "type": "message",
            "from": src,
            "to": dst,
            "arrow": arrow,
            "suffix": suffix,
            "text": label,
            "centralFrom": central_from,
            "centralTo": central_to,
        }

    return None


def ensure_participant(
    model: dict[str, Any],
    pid: str,
    *,
    label: str | None = None,
    kind: str = "participant",
    force_label: bool = False,
    force_kind: bool = False,
) -> None:
    participants = model["participants"]
    order = model["participant_order"]
    if pid not in participants:
        participants[pid] = {
            "id": pid,
            "label": label or pid,
            "kind": kind,
            "links": [],
            "properties": [],
            "details": [],
        }
        order.append(pid)
        return

    if label and force_label:
        participants[pid]["label"] = label
    elif label and participants[pid].get("label", pid) == pid:
        participants[pid]["label"] = label

    if kind and force_kind:
        participants[pid]["kind"] = kind
    elif kind and participants[pid].get("kind") == "participant":
        participants[pid]["kind"] = kind


def parse_box_header(line: str) -> tuple[str | None, str | None]:
    rest = line[3:].strip()
    if not rest:
        return None, None

    parts = rest.split(None, 1)
    first = parts[0]
    second = parts[1] if len(parts) > 1 else ""

    is_color = bool(re.fullmatch(r"#?[0-9a-fA-F]{6}", first)) or first.lower().startswith("rgb") or first.lower() in {
        "transparent",
        "white",
        "black",
        "gray",
        "grey",
        "red",
        "green",
        "blue",
        "yellow",
        "orange",
        "purple",
        "aqua",
        "teal",
    }

    if is_color:
        return first, second.strip() or None
    return None, rest


def parse_sequence_diagram(source: str) -> dict[str, Any]:
    model: dict[str, Any] = {
        "participants": {},
        "participant_order": [],
        "boxes": [],
        "events": [],
        "autonumber": {
            "enabled": False,
            "start": 1,
            "step": 1,
        },
        "meta": {
            "source": source,
            "title": "",
            "accTitle": "",
            "accDescr": "",
        },
    }

    # Mermaid allows semicolon-separated statements; treat ';' like a line break.
    lines = source.replace("\r\n", "\n").replace(";", "\n").split("\n")
    header_seen = False
    current_box: dict[str, Any] | None = None
    pending_create: list[str] = []
    pending_destroy: list[str] = []
    acc_descr_multiline = False
    acc_descr_lines: list[str] = []

    for raw_line in lines:
        line = raw_line.strip()
        if acc_descr_multiline:
            end_idx = line.find("}")
            if end_idx >= 0:
                prefix = line[:end_idx].strip()
                if prefix:
                    acc_descr_lines.append(prefix)
                model["meta"]["accDescr"] = clean_text("\n".join(acc_descr_lines))
                acc_descr_lines = []
                acc_descr_multiline = False
                line = line[end_idx + 1 :].strip()
            else:
                if line:
                    acc_descr_lines.append(line)
                continue

        if not line:
            continue
        if line.startswith("%%") or line.startswith("#"):
            continue

        if not header_seen:
            if re.match(r"^sequenceDiagram\b", line, flags=re.IGNORECASE):
                header_seen = True
                continue
            raise ValueError("sequenceDiagram header is required for sequence renderer")

        lower = line.lower()

        title_match = re.match(r"^title\s*:?\s*(.+)$", line, flags=re.IGNORECASE)
        if title_match:
            model["meta"]["title"] = clean_text(title_match.group(1))
            continue

        acc_title_match = re.match(r"^accTitle\s*:\s*(.+)$", line, flags=re.IGNORECASE)
        if acc_title_match:
            model["meta"]["accTitle"] = clean_text(acc_title_match.group(1))
            continue

        acc_descr_match = re.match(r"^accDescr\s*:\s*(.+)$", line, flags=re.IGNORECASE)
        if acc_descr_match:
            model["meta"]["accDescr"] = clean_text(acc_descr_match.group(1))
            continue

        acc_descr_start = re.match(r"^accDescr\s*\{\s*(.*)$", line, flags=re.IGNORECASE)
        if acc_descr_start:
            tail = acc_descr_start.group(1)
            if "}" in tail:
                model["meta"]["accDescr"] = clean_text(tail.split("}", 1)[0])
            else:
                acc_descr_multiline = True
                if tail.strip():
                    acc_descr_lines.append(tail.strip())
            continue

        if lower.startswith("autonumber") or lower == "sequencenumbers":
            if lower.startswith("autonumber"):
                parts = line.split()
                if len(parts) >= 2 and parts[1].lower() == "off":
                    model["autonumber"]["enabled"] = False
                    continue
                model["autonumber"]["enabled"] = True
                if len(parts) >= 2:
                    try:
                        model["autonumber"]["start"] = int(parts[1])
                    except Exception:
                        pass
                if len(parts) >= 3:
                    try:
                        model["autonumber"]["step"] = int(parts[2])
                    except Exception:
                        pass
            else:
                model["autonumber"]["enabled"] = True
            continue

        if lower.startswith("box"):
            color, label = parse_box_header(line)
            current_box = {
                "color": color,
                "label": label,
                "participants": [],
            }
            model["boxes"].append(current_box)
            continue

        if lower == "end" and current_box is not None:
            current_box = None
            continue

        participant_decl = parse_participant_decl(line)
        if participant_decl:
            kind, pid, label = participant_decl
            ensure_participant(model, pid, label=label, kind=kind, force_label=True, force_kind=True)
            if current_box is not None and pid not in current_box["participants"]:
                current_box["participants"].append(pid)
            continue

        create_decl = parse_create_decl(line)
        if create_decl:
            kind, pid, label = create_decl
            ensure_participant(model, pid, label=label, kind=kind, force_label=True, force_kind=True)
            if current_box is not None and pid not in current_box["participants"]:
                current_box["participants"].append(pid)
            pending_create.append(pid)
            continue

        destroy_decl = parse_destroy_decl(line)
        if destroy_decl:
            kind, pid, label = destroy_decl
            ensure_participant(model, pid, label=label, kind=kind)
            pending_destroy.append(pid)
            continue

        activate_match = re.match(r"^activate\s+(.+)$", line, flags=re.IGNORECASE)
        if activate_match:
            pid = parse_actor_ref(activate_match.group(1))
            ensure_participant(model, pid)
            model["events"].append({"type": "activate", "actor": pid})
            continue

        deactivate_match = re.match(r"^deactivate\s+(.+)$", line, flags=re.IGNORECASE)
        if deactivate_match:
            pid = parse_actor_ref(deactivate_match.group(1))
            ensure_participant(model, pid)
            model["events"].append({"type": "deactivate", "actor": pid})
            continue

        note_match = re.match(r"^note\s+(left of|right of|over)\s+([^:]+)\s*:\s*(.+)$", line, flags=re.IGNORECASE)
        if note_match:
            pos = note_match.group(1).lower()
            targets = [parse_actor_ref(part.strip()) for part in note_match.group(2).split(",") if part.strip()]
            for pid in targets:
                ensure_participant(model, pid)
            model["events"].append(
                {
                    "type": "note",
                    "position": pos,
                    "actors": targets,
                    "text": clean_text(note_match.group(3)),
                }
            )
            continue

        link_match = re.match(r"^link\s+([^:]+):\s*(.+?)\s*@\s*(\S+)\s*$", line, flags=re.IGNORECASE)
        if link_match:
            pid = parse_actor_ref(link_match.group(1).strip())
            ensure_participant(model, pid)
            model["participants"][pid]["links"].append({"label": clean_text(link_match.group(2)), "url": link_match.group(3)})
            continue

        links_match = re.match(r"^links\s+([^:]+):\s*(.+)$", line, flags=re.IGNORECASE)
        if links_match:
            pid = parse_actor_ref(links_match.group(1).strip())
            ensure_participant(model, pid)
            payload = links_match.group(2).strip()
            try:
                parsed = json.loads(payload)
                if isinstance(parsed, dict):
                    for k, v in parsed.items():
                        model["participants"][pid]["links"].append({"label": str(k), "url": str(v)})
            except Exception:
                pass
            continue

        prop_match = re.match(r"^(properties|details)\s+([^:]+)\s*:\s*(.+)$", line, flags=re.IGNORECASE)
        if prop_match:
            kind = prop_match.group(1).lower()
            pid = parse_actor_ref(prop_match.group(2).strip())
            ensure_participant(model, pid)
            payload = clean_text(prop_match.group(3))
            model["participants"][pid][kind].append(payload)
            continue

        block_start_match = re.match(r"^(loop|alt|opt|par|par_over|critical|break|rect)\b\s*(.*)$", line, flags=re.IGNORECASE)
        if block_start_match:
            kind = block_start_match.group(1).lower()
            label = clean_text(block_start_match.group(2) or "")
            event = {
                "type": "block_start",
                "kind": kind,
                "label": label,
            }
            if kind == "rect":
                color_match = re.match(r"^([#a-zA-Z0-9(),.]+)(?:\s+(.*))?$", label)
                if color_match:
                    event["color"] = color_match.group(1)
                    event["label"] = clean_text(color_match.group(2) or "")
            model["events"].append(event)
            continue

        block_branch_match = re.match(r"^(else|and|option)\b\s*(.*)$", line, flags=re.IGNORECASE)
        if block_branch_match:
            model["events"].append(
                {
                    "type": "block_branch",
                    "kind": block_branch_match.group(1).lower(),
                    "label": clean_text(block_branch_match.group(2) or ""),
                }
            )
            continue

        if lower == "end":
            model["events"].append({"type": "block_end"})
            continue

        message = parse_message(line)
        if message:
            ensure_participant(model, message["from"])
            ensure_participant(model, message["to"])
            message["create"] = pending_create.copy()
            message["destroy"] = pending_destroy.copy()
            pending_create.clear()
            pending_destroy.clear()
            model["events"].append(message)
            continue

    if acc_descr_multiline and acc_descr_lines:
        model["meta"]["accDescr"] = clean_text("\n".join(acc_descr_lines))

    return model


def parse_slide_size(slide_size: str | None) -> tuple[float, float]:
    text = (slide_size or "16:9").strip().lower()
    if text in {"16:9", "wide", "widescreen"}:
        return 13.333, 7.5
    if text in {"4:3", "standard"}:
        return 10.0, 7.5

    custom = re.fullmatch(r"([0-9]+(?:\.[0-9]+)?)x([0-9]+(?:\.[0-9]+)?)", text)
    if custom:
        return float(custom.group(1)), float(custom.group(2))

    return 13.333, 7.5


def choose_participant_shape(kind: str) -> int:
    key = (kind or "participant").lower()
    if key == "actor":
        return MSO_SHAPE.ROUNDED_RECTANGLE
    if key == "boundary":
        return MSO_SHAPE.CLOUD
    if key == "control":
        return MSO_SHAPE.HEXAGON
    if key == "entity":
        return MSO_SHAPE.RECTANGLE
    if key == "database":
        return MSO_SHAPE.CAN
    if key == "collections":
        return MSO_SHAPE.FLOWCHART_STORED_DATA
    if key == "queue":
        return MSO_SHAPE.FLOWCHART_DELAY
    return MSO_SHAPE.RECTANGLE


def apply_connector_style(connector: Any, *, color: str, dotted: bool, width_pt: float) -> None:
    connector.line.width = Pt(width_pt)
    connector.line.color.rgb = to_rgb(color)
    connector.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT if dotted else MSO_LINE_DASH_STYLE.SOLID


def draw_cross_marker(slide: Any, x: float, y: float, *, color: str, size: float = 0.06) -> None:
    for p0, p1 in [((x - size, y - size), (x + size, y + size)), ((x - size, y + size), (x + size, y - size))]:
        segment = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(p0[0]),
            Inches(p0[1]),
            Inches(p1[0]),
            Inches(p1[1]),
        )
        apply_connector_style(segment, color=color, dotted=False, width_pt=1.2)


def draw_central_marker(slide: Any, x: float, y: float, *, color: str) -> None:
    diameter = 0.07
    marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - diameter / 2.0), Inches(y - diameter / 2.0), Inches(diameter), Inches(diameter))
    marker.fill.solid()
    marker.fill.fore_color.rgb = to_rgb("FFFFFF")
    marker.line.color.rgb = to_rgb(color)
    marker.line.width = Pt(1.0)


def draw_actor_header(
    slide: Any,
    *,
    center_x: float,
    header_y: float,
    header_w: float,
    header_h: float,
    label: str,
    font_family: str,
    font_size: float,
) -> None:
    color = "334155"
    group = slide.shapes.add_group_shape()

    head_r = min(0.08, header_h * 0.18)
    head_cy = header_y + 0.10 + head_r
    body_top = head_cy + head_r + 0.01
    body_bottom = body_top + 0.13
    arm_y = body_top + 0.05
    leg_y = body_bottom + 0.09
    spread = 0.08

    head = group.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - head_r),
        Inches(head_cy - head_r),
        Inches(head_r * 2.0),
        Inches(head_r * 2.0),
    )
    head.fill.background()
    head.line.color.rgb = to_rgb(color)
    head.line.width = Pt(1.1)

    for p0, p1 in [
        ((center_x, body_top), (center_x, body_bottom)),
        ((center_x - spread, arm_y), (center_x + spread, arm_y)),
        ((center_x, body_bottom), (center_x - spread * 0.85, leg_y)),
        ((center_x, body_bottom), (center_x + spread * 0.85, leg_y)),
    ]:
        seg = group.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(p0[0]),
            Inches(p0[1]),
            Inches(p1[0]),
            Inches(p1[1]),
        )
        apply_connector_style(seg, color=color, dotted=False, width_pt=1.1)

    text_y = header_y + header_h * 0.58
    text_h = max(0.14, header_h * 0.40)
    text_box = group.shapes.add_textbox(Inches(center_x - header_w / 2.0), Inches(text_y), Inches(header_w), Inches(text_h))
    text_box.fill.background()
    text_box.line.fill.background()
    tf = text_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.color.rgb = to_rgb("0F172A")


def message_arrow_markers(arrow: str) -> tuple[str, str, bool, bool, bool]:
    dotted = arrow in DOTTED_ARROWS
    start = "none"
    end = "none"
    end_cross = False
    open_end = False

    if arrow in {"->>", "-->>"}:
        end = "triangle"
    elif arrow in {"<<->>", "<<-->>"}:
        start = "triangle"
        end = "triangle"
    elif arrow in {"-x", "--x"}:
        end_cross = True
    elif arrow in {"-)", "--)"}:
        open_end = True
    elif arrow in HALF_ARROW_FORWARD:
        end = "triangle"
    elif arrow in HALF_ARROW_REVERSE:
        start = "triangle"

    return start, end, dotted, end_cross, open_end


def estimate_note_size(text: str, font_size: float = 10.0) -> tuple[float, float]:
    lines = split_lines(text)
    if not lines:
        lines = [""]
    longest = max(text_units(line) for line in lines)
    w = clampf((longest * font_size * 0.64) / 72.0 + 0.24, 0.9, 3.0)
    h = clampf((len(lines) * font_size * 1.35) / 72.0 + 0.12, 0.28, 1.6)
    return w, h


def estimate_wrapped_lines(text: str, width_in: float, font_size_pt: float) -> int:
    lines = split_lines(text)
    if not lines:
        lines = [""]
    units_per_line = max(4.0, (max(0.25, width_in) * 72.0) / max(font_size_pt * 0.62, 1e-6))
    wrapped = 0
    for line in lines:
        units = max(1.0, text_units(line))
        wrapped += max(1, int(math.ceil(units / units_per_line)))
    return max(1, wrapped)


def estimate_message_label_size(text: str, font_size: float, max_width: float) -> tuple[float, float]:
    lines = split_lines(text)
    if not lines:
        lines = [""]
    longest = max(text_units(line) for line in lines)
    pref_w = (longest * font_size * 0.62) / 72.0 + 0.24
    width = clampf(pref_w, 0.55, max(0.70, max_width))
    wrapped = estimate_wrapped_lines(text, width, font_size)
    height = clampf((wrapped * font_size * 1.30) / 72.0 + 0.10, 0.18, 1.80)
    return width, height


def message_display_text(ev: dict[str, Any]) -> str:
    text = str(ev.get("text") or "")
    if ev.get("seqNo") is not None:
        return f"{ev['seqNo']}. {text}" if text else f"{ev['seqNo']}."
    return text


def render_sequence(
    model: dict[str, Any],
    output: Path,
    *,
    patch_text: str | None,
    slide_size: str | None,
    edge_routing: str | None,
    append_to: Path | None = None,
) -> None:
    append_target = append_to.resolve() if append_to else None
    if append_target and append_target.exists():
        prs = Presentation(str(append_target))
        slide_w = float(prs.slide_width) / EMU_PER_INCH
        slide_h = float(prs.slide_height) / EMU_PER_INCH
    else:
        prs = Presentation()
        slide_w, slide_h = parse_slide_size(slide_size)
        prs.slide_width = Inches(slide_w)
        prs.slide_height = Inches(slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    routing_mode = (edge_routing or "straight").strip().lower()
    if routing_mode not in {"straight", "elbow"}:
        routing_mode = "straight"

    font_family = "Yu Gothic UI"
    left_margin = 0.35
    right_margin = 0.35
    header_h = 0.52

    title_offset = 0.0
    title_text = clean_text(str(model.get("meta", {}).get("title") or ""))
    if title_text:
        title_box_h = 0.28
        title_shape = slide.shapes.add_textbox(Inches(0.35), Inches(0.10), Inches(max(0.8, slide_w - 0.70)), Inches(title_box_h))
        title_shape.fill.background()
        title_shape.line.fill.background()
        tf = title_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name = font_family
        run.font.size = Pt(13.5)
        run.font.bold = True
        run.font.color.rgb = to_rgb("0F172A")
        title_offset = 0.35

    header_y = 0.40 + title_offset
    top_margin = 0.35 + title_offset

    participant_ids = list(model["participant_order"])
    if not participant_ids:
        payload = "\n".join(["__MMD2PPTX_EMBED_START__", json.dumps({"sourceMmd": model["meta"]["source"], "patchYaml": patch_text, "version": 1}, ensure_ascii=False, indent=2), "__MMD2PPTX_EMBED_END__"])
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = payload
        save_path = append_target if append_target else output
        save_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(save_path))
        return

    n = len(participant_ids)
    available_w = max(1.0, slide_w - left_margin - right_margin)
    gap = 0.65
    header_w = (available_w - gap * (n - 1)) / max(1, n)
    if header_w < 0.82:
        gap = max(0.30, (available_w - 0.82 * n) / max(1, n - 1)) if n > 1 else 0.30
        header_w = max(0.82, (available_w - gap * (n - 1)) / max(1, n))
    header_w = min(2.1, header_w)

    total_used = header_w * n + gap * (n - 1)
    x0 = (slide_w - total_used) / 2.0

    participant_geom: dict[str, dict[str, float]] = {}
    for i, pid in enumerate(participant_ids):
        center = x0 + i * (header_w + gap) + header_w / 2.0
        participant_geom[pid] = {
            "center": center,
            "left": center - header_w / 2.0,
            "right": center + header_w / 2.0,
        }

    events = model["events"]
    timeline_top = header_y + header_h + 0.33
    message_font_pt = 8.8
    note_font_pt = 9.5
    spacing_scale = 1.0

    autonumber = model.get("autonumber", {})
    seq_enabled = bool(autonumber.get("enabled", False))
    seq_start = int(autonumber.get("start", 1))
    seq_step = int(autonumber.get("step", 1) or 1)

    def layout_timeline(
        *,
        message_font: float,
        note_font: float,
        spacing: float,
    ) -> tuple[float, list[dict[str, Any]], dict[str, float], dict[str, float]]:
        y = timeline_top
        blocks_local: list[dict[str, Any]] = []
        block_stack_local: list[dict[str, Any]] = []
        seq_current_local = seq_start
        create_y: dict[str, float] = {}
        destroy_y: dict[str, float] = {}

        for ev in events:
            ev_type = ev.get("type")
            ev["y"] = y

            if ev_type == "message":
                if seq_enabled:
                    ev["seqNo"] = seq_current_local
                    seq_current_local += seq_step
                else:
                    ev["seqNo"] = None

                from_id = ev["from"]
                to_id = ev["to"]
                for ctx in block_stack_local:
                    ctx["participants"].add(from_id)
                    ctx["participants"].add(to_id)

                for pid in ev.get("create", []):
                    create_y.setdefault(pid, y)
                for pid in ev.get("destroy", []):
                    destroy_y[pid] = y

                display_text = message_display_text(ev)
                ev["displayText"] = display_text
                if from_id == to_id:
                    max_label_w = clampf(header_w * 0.95, 0.72, 2.4)
                else:
                    span = abs(participant_geom[from_id]["center"] - participant_geom[to_id]["center"])
                    max_label_w = clampf(max(0.90, span - 0.18), 0.90, max(1.8, slide_w * 0.40))

                if display_text:
                    label_w, label_h = estimate_message_label_size(display_text, message_font, max_label_w)
                else:
                    label_w, label_h = 0.0, 0.0
                ev["labelW"] = label_w
                ev["labelH"] = label_h
                ev["labelFont"] = message_font

                base_step = (0.52 if from_id == to_id else 0.36) * spacing
                label_step = label_h + ((0.30 if from_id == to_id else 0.24) * spacing) if display_text else 0.0
                y += max(base_step, label_step)
                continue

            if ev_type == "note":
                for actor in ev.get("actors", []):
                    for ctx in block_stack_local:
                        ctx["participants"].add(actor)
                note_w, note_h = estimate_note_size(ev.get("text", ""), font_size=note_font)
                ev["noteW"] = note_w
                ev["noteH"] = note_h
                ev["noteFont"] = note_font
                y += max(0.28 * spacing, note_h + 0.08 * spacing)
                continue

            if ev_type in {"activate", "deactivate"}:
                actor = ev.get("actor")
                if actor:
                    for ctx in block_stack_local:
                        ctx["participants"].add(actor)
                y += 0.16 * spacing
                continue

            if ev_type == "block_start":
                ctx = {
                    "kind": ev.get("kind", "block"),
                    "label": ev.get("label", ""),
                    "color": ev.get("color"),
                    "startY": y - 0.05 * spacing,
                    "participants": set(),
                    "branches": [],
                    "depth": len(block_stack_local),
                }
                block_stack_local.append(ctx)
                y += 0.16 * spacing
                continue

            if ev_type == "block_branch":
                if block_stack_local:
                    block_stack_local[-1]["branches"].append(
                        {"y": y - 0.02 * spacing, "kind": ev.get("kind", ""), "label": ev.get("label", "")}
                    )
                y += 0.16 * spacing
                continue

            if ev_type == "block_end":
                if block_stack_local:
                    ctx = block_stack_local.pop()
                    ctx["endY"] = y + 0.05 * spacing
                    blocks_local.append(ctx)
                y += 0.12 * spacing
                continue

            y += 0.18 * spacing

        while block_stack_local:
            ctx = block_stack_local.pop()
            ctx["endY"] = y + 0.05 * spacing
            blocks_local.append(ctx)

        return y, blocks_local, create_y, destroy_y

    y_end, blocks, participant_create_y, participant_destroy_y = layout_timeline(
        message_font=message_font_pt,
        note_font=note_font_pt,
        spacing=spacing_scale,
    )

    available_bottom = slide_h - 0.35
    available_height = max(1.0, available_bottom - timeline_top)
    required_height = max(0.01, y_end - timeline_top)

    if required_height > available_height:
        ratio = available_height / required_height
        message_font_pt = max(5.8, message_font_pt * max(0.62, ratio))
        note_font_pt = max(6.2, note_font_pt * max(0.62, ratio))
        spacing_scale = max(0.62, ratio)
        y_end, blocks, participant_create_y, participant_destroy_y = layout_timeline(
            message_font=message_font_pt,
            note_font=note_font_pt,
            spacing=spacing_scale,
        )
        required_height = max(0.01, y_end - timeline_top)

    if required_height > available_height:
        y_scale = max(0.55, available_height / required_height)

        def scale_event_y(value: float) -> float:
            return timeline_top + (float(value) - timeline_top) * y_scale

        for ev in events:
            ev["y"] = scale_event_y(float(ev.get("y", timeline_top)))
            if ev.get("type") == "message":
                ev["labelH"] = max(0.14, float(ev.get("labelH", 0.20)) * y_scale)
                ev["labelFont"] = max(5.4, float(ev.get("labelFont", message_font_pt)) * y_scale)
            elif ev.get("type") == "note":
                ev["noteH"] = max(0.20, float(ev.get("noteH", 0.30)) * y_scale)
                ev["noteFont"] = max(6.0, float(ev.get("noteFont", note_font_pt)) * y_scale)

        for block in blocks:
            block["startY"] = scale_event_y(float(block.get("startY", timeline_top)))
            block["endY"] = scale_event_y(float(block.get("endY", timeline_top + 0.30)))
            for br in block.get("branches", []):
                br["y"] = scale_event_y(float(br.get("y", timeline_top)))

        participant_create_y = {pid: scale_event_y(v) for pid, v in participant_create_y.items()}
        participant_destroy_y = {pid: scale_event_y(v) for pid, v in participant_destroy_y.items()}
        y_end = timeline_top + required_height * y_scale

    diagram_bottom = min(slide_h - 0.35, y_end + 0.18)
    message_font_ratio = message_font_pt / 8.8
    header_font_pt = clampf(10.0 * message_font_ratio, 6.4, 10.0)
    block_font_pt = clampf(8.5 * message_font_ratio, 6.0, 8.5)
    branch_font_pt = clampf(8.0 * message_font_ratio, 5.8, 8.0)

    # Participant group boxes.
    for box in model.get("boxes", []):
        members = [pid for pid in box.get("participants", []) if pid in participant_geom]
        if not members:
            continue
        left = min(participant_geom[pid]["left"] for pid in members) - 0.12
        right = max(participant_geom[pid]["right"] for pid in members) + 0.12

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top_margin),
            Inches(max(0.15, right - left)),
            Inches(max(0.15, diagram_bottom - top_margin + 0.06)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = to_rgb(parse_css_color(str(box.get("color") or "EDF2F7"), "EDF2F7"))
        shape.fill.transparency = 88.0
        shape.line.color.rgb = to_rgb("94A3B8")
        shape.line.width = Pt(0.9)
        if box.get("label"):
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(box.get("label"))
            run.font.name = font_family
            run.font.size = Pt(block_font_pt)
            run.font.bold = True
            run.font.color.rgb = to_rgb("0F172A")

    # Draw blocks before messages.
    full_left = min(participant_geom[pid]["left"] for pid in participant_ids) - 0.16
    full_right = max(participant_geom[pid]["right"] for pid in participant_ids) + 0.16

    for block in sorted(blocks, key=lambda b: b.get("depth", 0)):
        ids = [pid for pid in block.get("participants", set()) if pid in participant_geom]
        if ids:
            left = min(participant_geom[pid]["left"] for pid in ids) - 0.14
            right = max(participant_geom[pid]["right"] for pid in ids) + 0.14
        else:
            left = full_left
            right = full_right

        y1 = clampf(float(block.get("startY", 1.0)), 0.25, slide_h - 0.25)
        y2 = clampf(float(block.get("endY", y1 + 0.3)), y1 + 0.10, slide_h - 0.20)

        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(y1),
            Inches(max(0.10, right - left)),
            Inches(max(0.10, y2 - y1)),
        )

        kind = str(block.get("kind", ""))
        if kind == "rect":
            rect.fill.solid()
            rect.fill.fore_color.rgb = to_rgb(parse_css_color(str(block.get("color") or "E2E8F0"), "E2E8F0"))
            rect.fill.transparency = 82.0
            rect.line.color.rgb = to_rgb("94A3B8")
            rect.line.width = Pt(0.8)
        else:
            rect.fill.background()
            rect.line.color.rgb = to_rgb("94A3B8")
            rect.line.width = Pt(0.9)
            rect.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        header = kind
        if block.get("label"):
            header = f"{kind} {block['label']}"
        if header.strip():
            tf = rect.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = header.strip()
            run.font.name = font_family
            run.font.size = Pt(block_font_pt)
            run.font.bold = True
            run.font.color.rgb = to_rgb("0F172A")

        for br in block.get("branches", []):
            yb = clampf(float(br.get("y", y1 + 0.12)), y1 + 0.08, y2 - 0.05)
            sep = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(left),
                Inches(yb),
                Inches(right),
                Inches(yb),
            )
            apply_connector_style(sep, color="94A3B8", dotted=True, width_pt=0.8)
            label = str(br.get("kind") or "").strip()
            if br.get("label"):
                label = f"{label} {br['label']}".strip()
            if label:
                tb = slide.shapes.add_textbox(Inches(left + 0.04), Inches(yb - 0.12), Inches(max(0.4, right - left - 0.08)), Inches(0.18))
                tb.fill.background()
                tb.line.fill.background()
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = False
                tf.auto_size = MSO_AUTO_SIZE.NONE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = label
                run.font.name = font_family
                run.font.size = Pt(branch_font_pt)
                run.font.color.rgb = to_rgb("334155")

    # Participants + lifelines.
    for pid in participant_ids:
        part = model["participants"][pid]
        geom = participant_geom[pid]
        label_text = str(part.get("label") or pid)
        kind = str(part.get("kind", "participant")).lower()
        if kind == "actor":
            draw_actor_header(
                slide,
                center_x=geom["center"],
                header_y=header_y,
                header_w=header_w,
                header_h=header_h,
                label=label_text,
                font_family=font_family,
                font_size=header_font_pt,
            )
        else:
            shape = slide.shapes.add_shape(
                choose_participant_shape(kind),
                Inches(geom["left"]),
                Inches(header_y),
                Inches(header_w),
                Inches(header_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = to_rgb("F8FAFC")
            shape.line.color.rgb = to_rgb("334155")
            shape.line.width = Pt(1.0)

            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label_text
            run.font.name = font_family
            run.font.size = Pt(header_font_pt)
            run.font.color.rgb = to_rgb("0F172A")

    lifeline_start_default = header_y + header_h + 0.08
    lifeline_end = diagram_bottom
    lifeline_start: dict[str, float] = {pid: lifeline_start_default for pid in participant_ids}

    lifeline_stop: dict[str, float] = {pid: lifeline_end for pid in participant_ids}
    for pid, dy in participant_destroy_y.items():
        if pid in lifeline_stop:
            lifeline_stop[pid] = min(lifeline_end, max(lifeline_start_default + 0.05, dy))

    for pid in participant_ids:
        x = participant_geom[pid]["center"]
        y1 = lifeline_start[pid]
        y2 = lifeline_stop[pid]
        if y2 <= y1:
            y2 = y1 + 0.01
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2))
        apply_connector_style(line, color="64748B", dotted=True, width_pt=1.0)

        if pid in participant_destroy_y:
            draw_cross_marker(slide, x, y2, color="64748B", size=0.05)

    # Event drawing + activation bars.
    activation_stack: dict[str, list[float]] = {pid: [] for pid in participant_ids}
    activation_depth: dict[str, int] = {pid: 0 for pid in participant_ids}

    def draw_activation(pid: str, y_start: float, y_end: float, depth_index: int) -> None:
        if y_end <= y_start + 0.01:
            return
        center = participant_geom[pid]["center"]
        width = 0.09
        x = center - width / 2.0 + depth_index * 0.055
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y_start), Inches(width), Inches(y_end - y_start))
        bar.fill.solid()
        bar.fill.fore_color.rgb = to_rgb("CFE8FF")
        bar.line.color.rgb = to_rgb("3B82F6")
        bar.line.width = Pt(0.8)

    for ev in events:
        ev_type = ev.get("type")
        y_ev = float(ev.get("y", 1.0))

        if ev_type == "activate":
            pid = ev.get("actor")
            if pid in activation_stack:
                activation_stack[pid].append(y_ev)
                activation_depth[pid] += 1
            continue

        if ev_type == "deactivate":
            pid = ev.get("actor")
            if pid in activation_stack and activation_stack[pid]:
                start_y = activation_stack[pid].pop()
                depth_index = max(0, len(activation_stack[pid]))
                draw_activation(pid, start_y, y_ev, depth_index)
                activation_depth[pid] = len(activation_stack[pid])
            continue

        if ev_type == "note":
            actors = [pid for pid in ev.get("actors", []) if pid in participant_geom]
            if not actors:
                continue
            text = str(ev.get("text") or "")
            note_w = float(ev.get("noteW", 0.0))
            note_h = float(ev.get("noteH", 0.0))
            note_font = float(ev.get("noteFont", note_font_pt))
            if note_w <= 0.0 or note_h <= 0.0:
                note_w, note_h = estimate_note_size(text, font_size=note_font)
            pos = str(ev.get("position") or "over").lower()

            if pos == "right of":
                anchor = participant_geom[actors[0]]["right"] + 0.12
                x = anchor
            elif pos == "left of":
                anchor = participant_geom[actors[0]]["left"] - 0.12
                x = anchor - note_w
            else:
                min_c = min(participant_geom[pid]["center"] for pid in actors)
                max_c = max(participant_geom[pid]["center"] for pid in actors)
                x = (min_c + max_c) / 2.0 - note_w / 2.0

            x = clampf(x, 0.08, slide_w - note_w - 0.08)
            y_note = clampf(y_ev - note_h / 2.0, 0.10, slide_h - note_h - 0.10)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y_note), Inches(note_w), Inches(note_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = to_rgb("FEF3C7")
            shape.line.color.rgb = to_rgb("D97706")
            shape.line.width = Pt(0.8)

            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.name = font_family
            run.font.size = Pt(note_font)
            run.font.color.rgb = to_rgb("7C2D12")
            continue

        if ev_type != "message":
            continue

        src = ev["from"]
        dst = ev["to"]
        if src not in participant_geom or dst not in participant_geom:
            continue

        x1 = participant_geom[src]["center"]
        x2 = participant_geom[dst]["center"]
        y_line = y_ev
        arrow = str(ev.get("arrow") or "->")
        start_marker, end_marker, dotted, end_cross, open_end = message_arrow_markers(arrow)
        if open_end:
            # "open" is not a valid OOXML line end type and causes PowerPoint repair.
            end_marker = "arrow"

        segments: list[Any] = []

        if src == dst:
            side_right = (slide_w - participant_geom[src]["right"]) >= participant_geom[src]["left"]
            loop_w = clampf(header_w * 0.65, 0.34, 0.92)
            loop_h = 0.24
            if side_right:
                pts = [
                    (x1, y_line),
                    (x1 + loop_w, y_line),
                    (x1 + loop_w, y_line + loop_h),
                    (x1, y_line + loop_h),
                ]
                label_x = x1 + loop_w * 0.72
            else:
                pts = [
                    (x1, y_line),
                    (x1 - loop_w, y_line),
                    (x1 - loop_w, y_line + loop_h),
                    (x1, y_line + loop_h),
                ]
                label_x = x1 - loop_w * 0.72

            for i in range(len(pts) - 1):
                p0, p1 = pts[i], pts[i + 1]
                seg = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(p0[0]),
                    Inches(p0[1]),
                    Inches(p1[0]),
                    Inches(p1[1]),
                )
                apply_connector_style(seg, color="1E293B", dotted=dotted, width_pt=1.3)
                segments.append(seg)

            if segments:
                if start_marker != "none":
                    set_line_markers(segments[0], start_marker, "none")
                set_line_markers(segments[-1], "none", end_marker if end_marker != "none" else "none")
                if end_cross:
                    draw_cross_marker(slide, pts[-1][0], pts[-1][1], color="1E293B", size=0.04)
                if ev.get("centralFrom"):
                    draw_central_marker(slide, pts[0][0], pts[0][1], color="1E293B")
                if ev.get("centralTo"):
                    draw_central_marker(slide, pts[-1][0], pts[-1][1], color="1E293B")

            text = str(ev.get("displayText") or message_display_text(ev))
            if text:
                tw = float(ev.get("labelW", 0.0))
                th = float(ev.get("labelH", 0.0))
                label_font = float(ev.get("labelFont", message_font_pt))
                if tw <= 0.0 or th <= 0.0:
                    tw, th = estimate_message_label_size(text, label_font, 2.2)
                tx = clampf(label_x - tw / 2.0, 0.08, slide_w - tw - 0.08)
                ty = clampf(y_line + 0.04, 0.08, slide_h - th - 0.08)
                tb = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
                tb.fill.background()
                tb.line.fill.background()
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = text
                run.font.name = font_family
                run.font.size = Pt(label_font)
                run.font.color.rgb = to_rgb("1E293B")

        else:
            # Sequence message lines are expected to stay horizontal.
            points: list[tuple[float, float]] = [(x1, y_line), (x2, y_line)]

            for i in range(len(points) - 1):
                p0, p1 = points[i], points[i + 1]
                seg = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(p0[0]),
                    Inches(p0[1]),
                    Inches(p1[0]),
                    Inches(p1[1]),
                )
                apply_connector_style(seg, color="1E293B", dotted=dotted, width_pt=1.25)
                if i == 0 and start_marker != "none":
                    set_line_markers(seg, start_marker, "none")
                segments.append(seg)

            if segments:
                last = segments[-1]
                set_line_markers(last, "none", end_marker if end_marker != "none" else "none")
                if end_cross:
                    draw_cross_marker(slide, points[-1][0], points[-1][1], color="1E293B", size=0.04)
                if ev.get("centralFrom"):
                    draw_central_marker(slide, points[0][0], points[0][1], color="1E293B")
                if ev.get("centralTo"):
                    draw_central_marker(slide, points[-1][0], points[-1][1], color="1E293B")

            text = str(ev.get("displayText") or message_display_text(ev))
            if text:
                tw = float(ev.get("labelW", 0.0))
                th = float(ev.get("labelH", 0.0))
                label_font = float(ev.get("labelFont", message_font_pt))
                if tw <= 0.0 or th <= 0.0:
                    max_label_w = clampf(max(0.90, abs(x2 - x1) - 0.18), 0.90, max(1.8, slide_w * 0.40))
                    tw, th = estimate_message_label_size(text, label_font, max_label_w)
                cx = (x1 + x2) / 2.0
                tx = clampf(cx - tw / 2.0, 0.08, slide_w - tw - 0.08)
                ty = clampf(y_line - th - 0.04, 0.08, slide_h - th - 0.08)
                tb = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
                tb.fill.background()
                tb.line.fill.background()
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = text
                run.font.name = font_family
                run.font.size = Pt(label_font)
                run.font.color.rgb = to_rgb("1E293B")

        # inline activation shortcut support.
        if ev.get("suffix") == "+":
            target = dst
            if target in activation_stack:
                activation_stack[target].append(y_ev + 0.02)
        elif ev.get("suffix") == "-":
            target = src
            if target in activation_stack and activation_stack[target]:
                start_y = activation_stack[target].pop()
                depth_index = max(0, len(activation_stack[target]))
                draw_activation(target, start_y, y_ev + 0.03, depth_index)

    for pid in participant_ids:
        while activation_stack[pid]:
            start_y = activation_stack[pid].pop()
            depth_index = max(0, len(activation_stack[pid]))
            draw_activation(pid, start_y, lifeline_stop.get(pid, diagram_bottom), depth_index)

    embed = {
        "sourceMmd": model["meta"].get("source", ""),
        "patchYaml": patch_text,
        "version": 1,
        "diagramType": "sequence",
    }
    payload = "\n".join(
        [
            "__MMD2PPTX_EMBED_START__",
            json.dumps(embed, ensure_ascii=False, indent=2),
            "__MMD2PPTX_EMBED_END__",
        ]
    )
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = payload

    save_path = append_target if append_target else output
    save_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(save_path))


def main() -> int:
    parser = argparse.ArgumentParser(description="Render Mermaid sequenceDiagram to PowerPoint using python-pptx")
    parser.add_argument("--source", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--patch", type=Path, default=None)
    parser.add_argument("--slide-size", type=str, default="16:9")
    parser.add_argument("--edge-routing", type=str, default="straight", choices=["straight", "elbow"])
    parser.add_argument("--append-to", type=Path, default=None)
    args = parser.parse_args()

    source_text = args.source.read_text(encoding="utf-8")
    patch_text = args.patch.read_text(encoding="utf-8") if args.patch else None
    model = parse_sequence_diagram(source_text)
    render_sequence(
        model,
        args.output,
        patch_text=patch_text,
        slide_size=args.slide_size,
        edge_routing=args.edge_routing,
        append_to=args.append_to,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
